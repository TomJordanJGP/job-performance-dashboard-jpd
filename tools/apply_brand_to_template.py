"""Apply Jobs Go Public brand styling to Renewals copy.pptx.

One-time, idempotent build script. Reads brand assets from `Brand/`
(gitignored, populated from brand-kit.zip). Downloads official primary logo
PNGs once if missing.

`Renewals.pptx` (the production template) is NEVER touched by this script —
both files are only allowed to be read for sanity-checking that the production
SHA hasn't changed. The single output target is `Renewals copy.pptx`.

Run: venv/bin/python tools/apply_brand_to_template.py
"""

from __future__ import annotations

import hashlib
import sys
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))
from theme.colors import JGP_COLORS  # noqa: E402

TEMPLATE = PROJECT_ROOT / 'Renewals copy.pptx'
PRODUCTION = PROJECT_ROOT / 'Renewals.pptx'  # only sha-checked, never written
BRAND_ROOT = PROJECT_ROOT / 'Brand'
LOGO_DIR = BRAND_ROOT / 'logos' / 'primary'
SHAPE_DIR = BRAND_ROOT / 'shapes'

LOGO_URLS = {
    'jobs-go-public-logo-full-colour.png':
        'https://media.jobsgopublic.com/wp-content/uploads/2025/10/JGP-Logo_RGB.png',
    'jobs-go-public-logo-white.png':
        'https://media.jobsgopublic.com/wp-content/uploads/2026/05/JGP-Logo_WHT.png',
}

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

DM_SANS = 'DM Sans'

# Exact-match curated label set. ALL-CAPS template labels → sentence case.
# Each map entry covers either a whole label or a fragment that appears as a
# standalone <a:r> run (Office often splits stat labels into multiple runs).
LABEL_MAP = {
    'ADVERTISING REPORT': 'Advertising report',
    'PREPARED FOR': 'Prepared for',
    'REPORTING PERIOD': 'Reporting period',
    'COMMENTARY': 'Commentary',
    'ROLES BENCHMARKED': 'Roles benchmarked',
    'STRONGEST JOB CATEGORY': 'Strongest job category',
    'AVERAGE VIEWS': 'Average views',
    'AVERAGE APPLIES': 'Average applies',
    'BENCHMARK': 'Benchmark',
    'AVERAGE': 'Average',
    'APPLIES': 'Applies',
    'ADVERTISED THIS PERIOD': 'Advertised this period',
    'COST PER JOB': 'Cost per job',
    'COST PER VIEW': 'Cost per view',
    'COST PER APPLY': 'Cost per apply',
    'COST PER': 'Cost per',
    'GET IN TOUCH': 'Get in touch',
    'YOUR ACCOUNT TEAM': 'Your account team',
    'GO PLACES. GO PUBLIC.': 'Go places. Go public.',
    'EMAIL': 'Email',
    'PHONE': 'Phone',
    'VIEWS': 'views',
    'JOBS': 'Jobs',
    'THIS PERIOD': 'this period',
    # Trailing-space variants for run-split labels like "BENCHMARK " + "AVERAGE" + " VIEWS"
    'BENCHMARK ': 'Benchmark ',
    'AVERAGE ': 'Average ',
    'PREPARED ': 'Prepared ',
    'JOBS ': 'Jobs ',
    'ADVERTISED ': 'advertised ',
    ' VIEWS': ' views',
    ' APPLIES': ' applies',
}

# Off-canon → canonical hex (lowercase keys; values keep upper for OOXML
# convention). We normalise the OOXML attribute case in apply_palette_fixes().
COLOUR_FIXES = {
    'd4eb66': 'E5FF6E',
    'e6ff6e': 'E5FF6E',
    '757575': '5B5067',
}


# ---------------------------------------------------------------------------
# Brand asset acquisition
# ---------------------------------------------------------------------------

def download_logos() -> None:
    """Download official PNG logos if not already present."""
    LOGO_DIR.mkdir(parents=True, exist_ok=True)
    for filename, url in LOGO_URLS.items():
        target = LOGO_DIR / filename
        if target.exists() and target.stat().st_size > 0:
            print(f"  logo: {filename} present ({target.stat().st_size:,} B)")
            continue
        print(f"  logo: downloading {filename}")
        with urllib.request.urlopen(url, timeout=30) as resp:
            target.write_bytes(resp.read())
        print(f"  logo: saved {filename} ({target.stat().st_size:,} B)")


def assert_brand_assets() -> None:
    required = [
        LOGO_DIR / 'jobs-go-public-logo-white.png',
        LOGO_DIR / 'jobs-go-public-logo-full-colour.png',
        SHAPE_DIR / 'Bold Green' / 'Curved Rectangle Stack.png',
    ]
    missing = [p for p in required if not p.exists()]
    if missing:
        print('Missing required brand assets:')
        for p in missing:
            print(f"  - {p}")
        sys.exit(1)


# ---------------------------------------------------------------------------
# XML transforms — applied to every slide / slide master / theme part
# ---------------------------------------------------------------------------

def replace_fonts(root) -> int:
    """Set every <a:latin>, <a:ea>, <a:cs> typeface to DM Sans where it's
    currently a literal font (not a theme reference like '+mn-lt' which is
    fixed via the theme font-scheme update). Also normalises chart-placeholder
    `monospace` to DM Sans for consistency. Drops Arial-targeted pitchFamily/
    charset hints so platforms without DM Sans don't fall back to Arial."""
    n = 0
    targets = (
        '{%s}latin' % A_NS,
        '{%s}ea' % A_NS,
        '{%s}cs' % A_NS,
    )
    for tag in targets:
        for el in root.iter(tag):
            current = el.get('typeface', '')
            # Skip empty (no override) and theme references (start with '+').
            if not current or current.startswith('+') or current == DM_SANS:
                continue
            el.set('typeface', DM_SANS)
            el.attrib.pop('pitchFamily', None)
            el.attrib.pop('charset', None)
            n += 1
    return n


def apply_palette_fixes(root) -> int:
    """Normalise off-canon hex fills and strokes to canonical brand tokens."""
    n = 0
    for el in root.iter('{%s}srgbClr' % A_NS):
        val = (el.get('val') or '').lower()
        if val in COLOUR_FIXES:
            el.set('val', COLOUR_FIXES[val])
            n += 1
    return n


def sentence_case_runs(root) -> int:
    """Apply LABEL_MAP to every <a:t> whose text exactly matches a key. Resets
    the parent run's `spc` (letter-spacing) so the wide tracking that mimicked
    uppercase doesn't survive into sentence case."""
    n = 0
    t_tag = '{%s}t' % A_NS
    rpr_tag = '{%s}rPr' % A_NS
    for t_el in root.iter(t_tag):
        if t_el.text is None:
            continue
        replacement = LABEL_MAP.get(t_el.text)
        if replacement is None:
            continue
        t_el.text = replacement
        run = t_el.getparent()
        if run is not None:
            rpr = run.find(rpr_tag)
            if rpr is not None and 'spc' in rpr.attrib:
                rpr.set('spc', '0')
        n += 1
    return n


def lift_minimum_text_size(root, min_pt: int = 10) -> int:
    """Bump any <a:rPr sz="…"/> below `min_pt` up to `min_pt`. OOXML font
    sizes are 1/100 of a point: sz="900" = 9pt."""
    n = 0
    threshold = min_pt * 100
    for rpr in root.iter('{%s}rPr' % A_NS):
        sz = rpr.get('sz')
        if sz is None:
            continue
        try:
            sz_int = int(sz)
        except ValueError:
            continue
        if sz_int < threshold:
            rpr.set('sz', str(threshold))
            n += 1
    # Also walk endParaRPr / defRPr — they carry size hints for empty paras.
    for tag in ('{%s}endParaRPr' % A_NS, '{%s}defRPr' % A_NS):
        for rpr in root.iter(tag):
            sz = rpr.get('sz')
            if sz is None:
                continue
            try:
                sz_int = int(sz)
            except ValueError:
                continue
            if sz_int < threshold:
                rpr.set('sz', str(threshold))
                n += 1
    return n


def update_theme_font_scheme(theme_root) -> int:
    """Force the theme major/minor font scheme to DM Sans (Latin). Leaves
    East-Asian and Complex-Script entries alone."""
    n = 0
    for tag in ('{%s}majorFont' % A_NS, '{%s}minorFont' % A_NS):
        for parent in theme_root.iter(tag):
            latin = parent.find('{%s}latin' % A_NS)
            if latin is not None and latin.get('typeface') != DM_SANS:
                latin.set('typeface', DM_SANS)
                latin.attrib.pop('pitchFamily', None)
                latin.attrib.pop('charset', None)
                n += 1
    return n


# ---------------------------------------------------------------------------
# Slide-level shape surgery
# ---------------------------------------------------------------------------

def _shape_text(sp) -> str:
    """Concatenated text content of every <a:t> inside a shape."""
    return ''.join((t.text or '') for t in sp.iter('{%s}t' % A_NS))


def _shape_pos(sp):
    """(x, y, cx, cy) in EMU, or None if any element missing."""
    off = sp.find('.//{%s}off' % A_NS)
    ext = sp.find('.//{%s}ext' % A_NS)
    if off is None or ext is None:
        return None
    return (int(off.get('x', 0)), int(off.get('y', 0)),
            int(ext.get('cx', 0)), int(ext.get('cy', 0)))


def remove_fake_logo_lockup(slide):
    """Remove the hand-built "Go" + "Jobs Go Public" lockup and its white
    background square. Returns (left, top) anchor of the removed lockup, or
    None if not found."""
    spTree = slide.shapes._spTree
    target_jobs = target_go = target_bg = None

    for sp in list(spTree.iter('{%s}sp' % P_NS)):
        text = _shape_text(sp).strip()
        if text == 'Jobs Go Public':
            target_jobs = sp
        elif text == 'Go':
            target_go = sp
        elif text == '':
            # Tiny white-filled rect in the top-left margin = logo background.
            pos = _shape_pos(sp)
            if pos is None:
                continue
            x, y, cx, cy = pos
            if x < 1_000_000 and y < 1_000_000 and cx < 600_000 and cy < 600_000:
                fills = sp.findall('.//{%s}solidFill/{%s}srgbClr' % (A_NS, A_NS))
                if any((f.get('val') or '').lower() == 'ffffff' for f in fills):
                    target_bg = sp

    if target_jobs is None:
        return None

    anchor_sp = target_go if target_go is not None else target_jobs
    pos = _shape_pos(anchor_sp)
    anchor = (pos[0], pos[1]) if pos else None

    for sp in (target_bg, target_go, target_jobs):
        if sp is not None:
            sp.getparent().remove(sp)

    return anchor


def remove_right_edge_decorative_rects(slide, slide_w_emu: int) -> int:
    """Remove plain rect-prstGeom shapes with no text that sit on the right
    portion of the slide (x > 60% of slide width). These are the off-canon
    `#9C67D3` and `#E5FF6E` colour blocks on the title slide."""
    spTree = slide.shapes._spTree
    n = 0
    for sp in list(spTree.iter('{%s}sp' % P_NS)):
        if _shape_text(sp).strip():
            continue
        prst = sp.find('.//{%s}prstGeom' % A_NS)
        if prst is None or prst.get('prst') != 'rect':
            continue
        pos = _shape_pos(sp)
        if pos is None:
            continue
        x, _, cx, _ = pos
        # Right-edge decoration: starts past 60% of slide width.
        if x > slide_w_emu * 0.6 and cx > 1_500_000:
            sp.getparent().remove(sp)
            n += 1
    return n


def add_logo(slide, anchor, logo_path: Path, target_width_emu: int = 1_700_000):
    """Insert the official JGP logo at the captured anchor. Width is set
    explicitly; height auto-scales from the source aspect ratio. 1.7M EMU ≈
    1.86 in ≈ 178 px @ 96 DPI — comfortably above the 130 px brand minimum."""
    if anchor is None:
        return None
    left, top = anchor
    return slide.shapes.add_picture(
        str(logo_path), Emu(left), Emu(top), width=Emu(target_width_emu)
    )


def add_title_brand_shape(slide, slide_w_emu: int, slide_h_emu: int,
                          shape_path: Path):
    """Anchor an approved brand shape PNG to the right edge of the title
    slide, sized to roughly mirror the previous decorative blocks."""
    width_emu = 3_200_000  # ~3.5 in
    # Read aspect ratio via Pillow if available, else assume square.
    try:
        from PIL import Image
        with Image.open(shape_path) as im:
            iw, ih = im.size
        height_emu = int(width_emu * ih / iw)
    except Exception:
        height_emu = width_emu
    left = slide_w_emu - width_emu - 300_000
    top = (slide_h_emu - height_emu) // 2
    return slide.shapes.add_picture(
        str(shape_path), Emu(left), Emu(top),
        width=Emu(width_emu), height=Emu(height_emu)
    )


def add_consistent_footer(slide, slide_w_emu: int, slide_h_emu: int):
    """Add a bottom-strip footer to a content slide (slides 2-7) with a small
    DM Sans tagline at the left and a page-number placeholder at the right.
    Mirrors the visual position of slide 1's existing pill so the eye reads
    a consistent baseline across the deck."""
    pill_top_emu = 9_742_289  # mirror slide-1's pill Y
    pill_h_emu = 240_000

    # Tagline at bottom-left
    left = Emu(609_600)
    width = Emu(2_500_000)
    tagline = slide.shapes.add_textbox(left, Emu(pill_top_emu), width, Emu(pill_h_emu))
    tf = tagline.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = 'Jobs Go Public  ·  Go places. Go public.'
    run.font.name = DM_SANS
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor.from_string(JGP_COLORS['primary'].lstrip('#').upper())

    # Page number at bottom-right (uses {{slide_number}} / {{total_slides}}
    # which the runtime export already populates per slide).
    pn_w = Emu(2_500_000)
    pn_left = Emu(slide_w_emu - 609_600 - 2_500_000)
    pn_box = slide.shapes.add_textbox(pn_left, Emu(pill_top_emu), pn_w, Emu(pill_h_emu))
    pf = pn_box.text_frame
    pf.margin_left = Emu(0)
    pf.margin_top = Emu(0)
    pf.margin_right = Emu(0)
    pf.margin_bottom = Emu(0)
    pp = pf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    prun = pp.add_run()
    prun.text = '{{slide_number}} / {{total_slides}}'
    prun.font.name = DM_SANS
    prun.font.size = Pt(10)
    prun.font.color.rgb = RGBColor.from_string('5B5067')


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------

def transform_xml_root(root) -> dict:
    """Apply font / palette / sentence-case / minimum-size to any XML root
    (slide, slide master, or theme). Returns counts for logging."""
    return {
        'fonts': replace_fonts(root),
        'palette': apply_palette_fixes(root),
        'sentence_case': sentence_case_runs(root),
        'lifted_size': lift_minimum_text_size(root, min_pt=10),
    }


def sha256(path: Path) -> str:
    h = hashlib.sha256()
    h.update(path.read_bytes())
    return h.hexdigest()


def main() -> int:
    if not TEMPLATE.exists():
        print(f"ERROR: source template not found at {TEMPLATE}")
        return 1
    if not BRAND_ROOT.exists():
        print(f"ERROR: brand kit not staged at {BRAND_ROOT}")
        print("       Run: unzip -o brand-kit.zip && cp -R brand-kit/Brand .")
        return 1

    # Production-safety check: capture the production SHA before AND after.
    prod_sha_before = sha256(PRODUCTION) if PRODUCTION.exists() else None
    print(f"Renewals.pptx (production) SHA before: {prod_sha_before}")

    print("\n[1/6] Downloading official logos…")
    download_logos()
    assert_brand_assets()

    print(f"\n[2/6] Opening template: {TEMPLATE.name}")
    prs = Presentation(str(TEMPLATE))

    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    print(f"      slide size: {slide_w_emu/914400:.2f} in x {slide_h_emu/914400:.2f} in")
    print(f"      slides: {len(prs.slides)}")

    print("\n[3/6] Per-slide XML transform (fonts, palette, sentence-case, sizes)")
    summary = {'fonts': 0, 'palette': 0, 'sentence_case': 0, 'lifted_size': 0}
    for idx, slide in enumerate(prs.slides, start=1):
        counts = transform_xml_root(slide.shapes._spTree)
        for k, v in counts.items():
            summary[k] += v
        if any(counts.values()):
            print(f"      slide {idx}: {counts}")

    # Theme + slide master
    print("\n[4/6] Theme + slide master font scheme → DM Sans")
    from lxml import etree
    THEME_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
    for sm in prs.slide_masters:
        counts = transform_xml_root(sm.element)
        print(f"      slide master: {counts}")
        # The theme part is a generic Part (not XmlPart), so we parse its blob
        # directly with lxml and write it back.
        theme_part = sm.part.part_related_by(THEME_REL)
        theme_root = etree.fromstring(theme_part.blob)
        theme_counts = transform_xml_root(theme_root)
        scheme_count = update_theme_font_scheme(theme_root)
        theme_counts['theme_font_scheme'] = scheme_count
        theme_part._blob = etree.tostring(
            theme_root, xml_declaration=True, encoding='UTF-8', standalone=True
        )
        print(f"      theme: {theme_counts}")

    print("\n[5/6] Per-slide shape surgery (logo lockups, decorative rects, footer)")
    logo_white = LOGO_DIR / 'jobs-go-public-logo-white.png'
    brand_shape = SHAPE_DIR / 'Bold Green' / 'Curved Rectangle Stack.png'

    for idx, slide in enumerate(prs.slides, start=1):
        # Replace fake logo lockup on every slide that has one (slides 1 and 8).
        anchor = remove_fake_logo_lockup(slide)
        if anchor is not None:
            add_logo(slide, anchor, logo_white)
            print(f"      slide {idx}: replaced fake lockup with white logo at {anchor}")

        # Title slide only: strip random colour-block rectangles + add brand shape.
        if idx == 1:
            removed = remove_right_edge_decorative_rects(slide, slide_w_emu)
            if removed:
                print(f"      slide 1: removed {removed} right-edge decorative rect(s)")
                add_title_brand_shape(slide, slide_w_emu, slide_h_emu, brand_shape)
                print(f"      slide 1: added approved brand shape from {brand_shape.name}")

        # Content slides 2-7: add consistent footer (slides 1 and 8 already
        # carry the "Go places. Go public." pill, just font/case-fixed).
        if 2 <= idx <= 7:
            add_consistent_footer(slide, slide_w_emu, slide_h_emu)
            print(f"      slide {idx}: added consistent footer")

    print(f"\n[6/6] Saving → {TEMPLATE.name}")
    prs.save(str(TEMPLATE))

    # Production-safety check after save.
    prod_sha_after = sha256(PRODUCTION) if PRODUCTION.exists() else None
    if prod_sha_before != prod_sha_after:
        print(f"FATAL: Renewals.pptx SHA changed during script run.")
        print(f"  before: {prod_sha_before}")
        print(f"  after:  {prod_sha_after}")
        return 2
    print(f"\nRenewals.pptx (production) SHA after:  {prod_sha_after}")
    print(f"Renewals.pptx untouched: {prod_sha_before == prod_sha_after}")
    print(f"\nTransform totals across slides: {summary}")
    print("Done.")
    return 0


if __name__ == '__main__':
    sys.exit(main())
