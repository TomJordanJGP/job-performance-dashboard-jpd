"""Client Report tab — branded PPTX advertising report for renewals.

Renders the five-section advertising report (Benchmarking Scatter,
Benchmarking Summary, Job Postings, Advertising ROI, Media Performance)
with optional PowerPoint export driven by the Renewals.pptx template.
"""

import io
import re
from datetime import datetime, timedelta

import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from plotly.subplots import make_subplots

from pptx import Presentation

from theme.colors import JGP_COLORS, JGP_PLOTLY_TEMPLATE
from theme.components import (
    client_hero,
    commentary_panel,
    export_cta_panel,
    kpi_card,
    kpi_card_dark,
    section_eyebrow,
    status_grid,
    summary_bar,
)
from data.processing import apply_media_categories
from data.loader import load_client_hq_regions


# Static explainers describing how each chart is calculated. Single source of
# truth: rendered on-screen as captions and substituted into the PPTX template
# via {{chart_explainer_<key>}} placeholders. Wording change → one-line edit
# here, no PowerPoint round trip.
CHART_EXPLAINERS = {
    'section_01_headlines': (
        "Figures that summarise the client's overall performance for the "
        "reporting period — volume, engagement, and cost efficiency at a glance."
    ),
    'section_02_intro': (
        "How each individual vacancy performed against the market average "
        "for the same occupation, plotted on views and applies and grouped "
        "into engagement categories."
    ),
    'section_03_intro': (
        "The client's average views and applies per vacancy compared with "
        "the wider market benchmark, summarised as headline indices and a "
        "side-by-side comparison chart."
    ),
    'section_04_intro': (
        "Volume of jobs posted alongside the apply clicks they generated, "
        "broken down by occupation category over the reporting period."
    ),
    'section_05_intro': (
        "Spend efficiency at a glance — total spend versus the rate-card "
        "equivalent, the cash saving delivered, and the cost paid per apply."
    ),
    'section_06_intro': (
        "Which occupations produced applies most efficiently and which ones "
        "cost the most per apply, with the full breakdown available on demand."
    ),
    'section_07_intro': (
        "Where the client's advertised salaries sit against the wider market "
        "for the most-posted occupations, with the client mean, national mean, "
        "and regional mean marked on each distribution."
    ),
    'section_08_intro': (
        "Performance of the traffic channels that fed candidates to the "
        "client's vacancies, ranked by conversion from views to applies."
    ),
    'section_09_intro': (
        "Download a branded PowerPoint deck containing every chart and "
        "commentary on this page, ready to share with the client."
    ),
    'benchmark_scatter': (
        "Each marker is one of your vacancies. Its position shows how views "
        "and applies compare to the average for the same occupation across "
        "all other clients — top-right is above benchmark on both."
    ),
    'benchmark_average': (
        "Your average views and applies per vacancy as a percentage of the "
        "wider market average. 100% sits in line with the benchmark; above "
        "100% outperforms it."
    ),
    'postings_by_type': (
        "How your vacancy volume and apply clicks distribute across "
        "occupation categories during the report period."
    ),
    'spend_vs_ratecard': (
        "Your subscription spend (purple) stacked with the saving versus "
        "paying rate-card per vacancy (green). The full bar is what these "
        "postings would cost without your subscription."
    ),
    'cost_per_app_by_occupation': (
        "Your annual spend allocated to each occupation by share of "
        "vacancies, then divided by the applies generated. Lower bars "
        "indicate where spend produces candidates most efficiently."
    ),
    'media_performance': (
        "Average views and applies per vacancy, broken down by traffic "
        "source. Shows which channels (organic search, paid, direct, "
        "referral, etc.) drive the most candidates."
    ),
    'salary_by_occupation': (
        "Market salary spread for your top 10 most-priced occupations. "
        "Lines mark your mean (red), the national mean (blue) and your "
        "HQ-region mean (green). Ranked by count of priced vacancies; "
        "minimum 5 per occupation."
    ),
}


@st.cache_data(ttl=3600)
def _read_chart_slot_dimensions(template_path: str, _mtime: float) -> dict:
    """Inner cached reader. The `_mtime` arg is part of the cache key so the
    cache invalidates automatically when the PPTX template is re-saved.
    """
    try:
        prs = Presentation(template_path)
    except Exception:
        return {}
    EMU_PER_CM = 360000
    EMU_PER_PX = 9525
    slots = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            m = re.search(r'\{\{chart:([^}]+)\}\}', text)
            if not m:
                continue
            slot = m.group(1)
            slots[slot] = {
                'cm': (shape.width / EMU_PER_CM, shape.height / EMU_PER_CM),
                'px': (int(shape.width / EMU_PER_PX), int(shape.height / EMU_PER_PX)),
            }
    return slots


def get_chart_slot_dimensions(template_path: str = 'Renewals.pptx') -> dict:
    """Return the on-slide size (cm + px) of every {{chart:slot}} placeholder.

    Re-reads the template whenever the file is modified (cache key includes
    mtime). Used to surface placeholder geometry on the dashboard so chart
    text sizing can be reasoned about against the actual rendered PPTX
    (text px in the export equals displayed px in the slide).

    Returns:
        {slot_name: {'cm': (w, h), 'px': (w, h)}} or {} if the template can't
        be opened.
    """
    import os
    try:
        mtime = os.path.getmtime(template_path)
    except OSError:
        mtime = 0.0
    return _read_chart_slot_dimensions(template_path, mtime)


def chart_caption(slot_name: str, slot_dims: dict) -> str:
    """Build the on-screen caption for a chart: explainer + PPTX slot size."""
    explainer = CHART_EXPLAINERS.get(slot_name, '')
    dim = slot_dims.get(slot_name)
    if not dim:
        return explainer
    cm_w, cm_h = dim['cm']
    px_w, px_h = dim['px']
    size_line = (f"_PPTX placeholder: {cm_w:.1f} × {cm_h:.1f} cm "
                 f"({px_w} × {px_h} px on slide)._")
    return f"{explainer}  \n{size_line}"


def generate_section_commentary(section, data):
    """Generate template-based commentary for each report section.

    Args:
        section: One of 'scatter', 'benchmark', 'postings', 'roi', 'media'.
        data: Dict with relevant metrics for the section.

    Returns:
        Markdown-formatted string with 2-4 sentences of data-driven insight.
    """
    fallback = "Insufficient data for detailed commentary."

    if section == 'scatter':
        total = data.get('total_count', 0)
        benchmarkable = data.get('benchmarkable_count', 0)
        zero_applies = data.get('zero_applies_count', 0)
        no_benchmark = data.get('no_benchmark_count', 0)
        top_performers = data.get('top_performers', [])
        worst_performers = data.get('worst_performers', [])

        if total == 0:
            return fallback
        if total == 1:
            return ("Only one vacancy is available for analysis, which limits "
                    "benchmarking comparisons. A larger sample of roles will "
                    "enable more meaningful performance insights.")

        parts = []
        parts.append(
            f"Of **{total}** vacancies, **{benchmarkable}** "
            f"({'all' if benchmarkable == total else f'{benchmarkable / total:.0%}'}) "
            f"have sufficient market data for benchmarking"
            f"{f', while **{no_benchmark}** lack benchmark data due to low sample sizes in their occupation category' if no_benchmark > 0 else ''}."
        )

        if top_performers:
            top = top_performers[0]
            parts.append(
                f"**{top['title']}** ({top['occupation']}) is a standout performer, "
                f"exceeding the benchmark by {top['views_diff_pct']:+.0f}% on views "
                f"and {top['applies_diff_pct']:+.0f}% on applies."
            )

        if zero_applies > 0:
            pct_zero = zero_applies / total * 100
            parts.append(
                f"**{zero_applies}** role{'s' if zero_applies != 1 else ''} "
                f"({pct_zero:.0f}% of total) received zero apply clicks and "
                f"may benefit from revised job descriptions or enhanced visibility."
            )

        return " ".join(parts)

    elif section == 'benchmark':
        client_clicks = data.get('client_avg_clicks', 0)
        bench_clicks = data.get('benchmark_avg_clicks', 0)
        client_applies = data.get('client_avg_applies', 0)
        bench_applies = data.get('benchmark_avg_applies', 0)
        num_jobs = data.get('num_jobs', 0)
        client_name = data.get('client_name', 'This client')

        if num_jobs == 0 or bench_clicks == 0:
            return fallback

        views_diff = ((client_clicks - bench_clicks) / bench_clicks) * 100
        views_word = "more" if views_diff >= 0 else "fewer"

        parts = []
        parts.append(
            f"Your vacancies received **{abs(views_diff):.0f}% {views_word} views** "
            f"than the market average, "
            f"{'indicating strong visibility across the platform' if views_diff >= 0 else 'suggesting opportunities to improve listing visibility'}."
        )

        if bench_applies > 0:
            applies_diff = ((client_applies - bench_applies) / bench_applies) * 100
            applies_word = "above" if applies_diff >= 0 else "below"
            parts.append(
                f"Apply rates sit **{abs(applies_diff):.0f}% {applies_word} benchmark**"
                f"{', reflecting strong candidate engagement' if applies_diff >= 10 else ', suggesting job descriptions may benefit from enhancement' if applies_diff < -10 else ', broadly in line with market expectations'}."
            )

        parts.append(
            f"This analysis is based on **{num_jobs}** vacancies posted by {client_name} "
            f"during the report period."
        )
        return " ".join(parts)

    elif section == 'postings':
        total_jobs = data.get('total_jobs', 0)
        total_applies = data.get('total_applies', 0)
        by_type = data.get('by_type')
        client_name = data.get('client_name', 'This client')

        if total_jobs == 0:
            return fallback
        if by_type is None or len(by_type) == 0:
            return f"{client_name} posted **{total_jobs}** vacancies with **{total_applies:,}** total apply clicks during this period."

        sorted_by_applies = by_type.sort_values('apply_clicks', ascending=False)
        top = sorted_by_applies.iloc[0]
        parts = []

        if len(by_type) == 1:
            parts.append(
                f"All **{total_jobs}** vacancies fall under **{top['occupation']}**, "
                f"generating **{int(top['apply_clicks']):,}** apply clicks."
            )
        else:
            parts.append(
                f"**{top['occupation']}** leads apply generation with "
                f"**{int(top['apply_clicks']):,}** apply clicks across "
                f"**{int(top['jobs_posted'])}** postings."
            )
            bottom = sorted_by_applies.iloc[-1]
            if int(bottom['apply_clicks']) == 0:
                parts.append(
                    f"**{bottom['occupation']}** received no apply clicks despite "
                    f"**{int(bottom['jobs_posted'])}** postings — these roles may "
                    f"benefit from revised job titles or enhanced descriptions."
                )
            elif len(by_type) > 1:
                parts.append(
                    f"**{bottom['occupation']}** generated the fewest applies "
                    f"(**{int(bottom['apply_clicks']):,}**) and may warrant "
                    f"targeted improvements to boost candidate engagement."
                )

        if total_applies == 0:
            parts.append(
                "No apply clicks were recorded across any category — reviewing "
                "listing quality and distribution channels is recommended."
            )
        return " ".join(parts)

    elif section == 'roi':
        annual_spend = data.get('annual_spend', 0)
        num_jobs = data.get('num_jobs', 0)
        cost_per_apply = data.get('cost_per_apply', 0)
        saving_pct = data.get('saving_pct', 0)
        roi_by_type = data.get('roi_by_type')

        if annual_spend == 0 or num_jobs == 0:
            return "Enter your annual spend and rate card price to generate ROI commentary."

        parts = []
        if saving_pct > 0:
            parts.append(
                f"Your advertising investment delivers a **{saving_pct:.0f}% saving** "
                f"compared to rate card pricing, demonstrating strong value from "
                f"the platform partnership."
            )
        else:
            parts.append(
                "Current spend exceeds rate card value — reviewing the pricing "
                "structure or consolidating lower-performing listings may improve "
                "overall return on investment."
            )

        total_applies = data.get('total_applies', 0)
        if total_applies > 0:
            parts.append(
                f"At **£{cost_per_apply:,.2f} per apply**, each candidate "
                f"enquiry represents a cost-effective acquisition channel."
            )

        if roi_by_type is not None and len(roi_by_type) > 1:
            best = roi_by_type.iloc[0]
            worst = roi_by_type.iloc[-1]
            parts.append(
                f"**{best['occupation']}** achieves the best cost efficiency at "
                f"£{best['cost_per_apply']:,.2f} per apply, while "
                f"**{worst['occupation']}** is the most expensive at "
                f"£{worst['cost_per_apply']:,.2f}."
            )

        return " ".join(parts)

    elif section == 'media':
        cat_stats = data.get('cat_stats')
        client_name = data.get('client_name', 'This client')

        if cat_stats is None or len(cat_stats) == 0:
            return "Media source data is not yet available for this client."

        sorted_stats = cat_stats.sort_values('total_applies', ascending=False)
        top_source = sorted_stats.iloc[0]
        parts = []

        parts.append(
            f"**{top_source['source_category']}** is the leading traffic source, "
            f"generating **{int(top_source['total_applies']):,}** apply clicks "
            f"from **{int(top_source['total_clicks']):,}** views."
        )

        best_conv = cat_stats.loc[cat_stats['conversion_rate'].idxmax()]
        if best_conv['conversion_rate'] > 0:
            parts.append(
                f"**{best_conv['source_category']}** achieves the highest "
                f"view-to-apply conversion rate at **{best_conv['conversion_rate']:.1f}%**."
            )

        paid_rows = cat_stats[cat_stats['source_category'].str.contains('Paid|PPC|Sponsored', case=False, na=False)]
        if len(paid_rows) > 0:
            paid = paid_rows.iloc[0]
            parts.append(
                f"Paid channels ({paid['source_category']}) contributed "
                f"**{int(paid['total_applies']):,}** applies with a "
                f"**{paid['conversion_rate']:.1f}%** conversion rate."
            )

        return " ".join(parts)

    return fallback


def generate_section_commentary_structured(section, data):
    """Structured commentary for PPTX template — returns dict with intro + bullet points.

    Returns: dict with keys 'intro', 'point_1', 'point_2', 'point_3' (last is optional).
    Each value is plain text (no markdown). Empty string for unused points.
    """
    def _clean(text):
        """Strip markdown bold markers."""
        return re.sub(r'\*\*(.+?)\*\*', r'\1', text or '')

    if section == 'benchmark_scatter':
        total = data.get('total_count', 0)
        benchmarkable = data.get('benchmarkable_count', 0)
        zero_applies = data.get('zero_applies_count', 0)
        no_benchmark = data.get('no_benchmark_count', 0)
        top_performers = data.get('top_performers', [])
        client_name = data.get('client_name', 'This client')

        if total == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': '', 'point_3': ''}

        intro = (f"{client_name} posted {total} vacancies during this period. "
                 f"Of these, {benchmarkable} have sufficient market data to benchmark against comparable public sector roles.")

        point_1 = ''
        if top_performers:
            top = top_performers[0]
            point_1 = (f"{top['title']} ({top['occupation']}) is the standout performer, "
                       f"exceeding the benchmark by {top['views_diff_pct']:+.0f}% on views and "
                       f"{top['applies_diff_pct']:+.0f}% on applies.")

        point_2 = ''
        if zero_applies > 0:
            pct = zero_applies / total * 100 if total > 0 else 0
            point_2 = (f"{zero_applies} role{'s' if zero_applies != 1 else ''} ({pct:.0f}% of total) "
                       f"received zero apply clicks — these may benefit from revised job descriptions or enhanced visibility.")

        point_3 = ''
        if no_benchmark > 0:
            point_3 = (f"{no_benchmark} role{'s' if no_benchmark != 1 else ''} could not be benchmarked "
                       f"due to low market sample sizes in their occupation category.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    elif section == 'benchmark_average':
        client_clicks = data.get('client_avg_clicks', 0)
        bench_clicks = data.get('benchmark_avg_clicks', 0)
        client_applies = data.get('client_avg_applies', 0)
        bench_applies = data.get('benchmark_avg_applies', 0)
        num_jobs = data.get('num_jobs', 0)
        client_name = data.get('client_name', 'This client')

        if num_jobs == 0 or bench_clicks == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': ''}

        views_diff = ((client_clicks - bench_clicks) / bench_clicks) * 100
        intro = (f"Across {num_jobs} vacancies, {client_name} averaged {client_clicks:,.0f} views "
                 f"and {client_applies:,.1f} applies per role.")

        views_word = "more" if views_diff >= 0 else "fewer"
        point_1 = (f"Your vacancies received {abs(views_diff):.0f}% {views_word} views than the market average — "
                   f"{'indicating strong visibility across the platform' if views_diff >= 0 else 'suggesting opportunities to improve listing visibility'}.")

        point_2 = ''
        if bench_applies > 0:
            applies_diff = ((client_applies - bench_applies) / bench_applies) * 100
            applies_word = "above" if applies_diff >= 0 else "below"
            tone = ('reflecting strong candidate engagement' if applies_diff >= 10
                    else 'suggesting job descriptions may benefit from enhancement' if applies_diff < -10
                    else 'broadly in line with market expectations')
            point_2 = f"Apply rates sit {abs(applies_diff):.0f}% {applies_word} benchmark, {tone}."

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'postings':
        total_jobs = data.get('total_jobs', 0)
        total_applies = data.get('total_applies', 0)
        by_type = data.get('by_type')
        client_name = data.get('client_name', 'This client')

        if total_jobs == 0 or by_type is None or len(by_type) == 0:
            return {'intro': 'Insufficient data for commentary.', 'point_1': '', 'point_2': ''}

        sorted_by_applies = by_type.sort_values('apply_clicks', ascending=False)
        top = sorted_by_applies.iloc[0]
        intro = (f"{client_name} posted {total_jobs} vacancies generating {total_applies:,} apply clicks across "
                 f"{len(by_type)} occupation categories.")

        point_1 = (f"{top['occupation']} leads apply generation with {int(top['apply_clicks']):,} apply clicks across "
                   f"{int(top['jobs_posted'])} postings — your strongest performing category.")

        point_2 = ''
        if len(by_type) > 1:
            bottom = sorted_by_applies.iloc[-1]
            if int(bottom['apply_clicks']) == 0:
                point_2 = (f"{bottom['occupation']} received no apply clicks despite {int(bottom['jobs_posted'])} postings — "
                           f"these roles may benefit from revised titles or enhanced descriptions.")
            else:
                point_2 = (f"{bottom['occupation']} generated the fewest applies ({int(bottom['apply_clicks']):,}) "
                           f"and may warrant targeted improvements.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'roi':
        annual_spend = data.get('annual_spend', 0)
        rate_card_total = data.get('rate_card_total', 0)
        num_jobs = data.get('num_jobs', 0)
        cost_per_apply = data.get('cost_per_apply', 0)
        saving_pct = data.get('saving_pct', 0)
        roi_by_type = data.get('roi_by_type')

        if annual_spend == 0 or num_jobs == 0:
            return {'intro': 'Enter your annual spend and rate card price to generate ROI commentary.',
                    'point_1': '', 'point_2': ''}

        saving_amount = rate_card_total - annual_spend
        intro = (f"Across {num_jobs} vacancies, your subscription delivered £{saving_amount:,.0f} of value "
                 f"compared to rate card pricing — a {saving_pct:.0f}% saving.")

        point_1 = (f"At £{cost_per_apply:,.2f} per apply, each candidate enquiry represents a "
                   f"cost-effective acquisition channel for {data.get('client_name', 'your team')}.")

        point_2 = ''
        if roi_by_type is not None and len(roi_by_type) > 1:
            best = roi_by_type.iloc[0]
            worst = roi_by_type.iloc[-1]
            point_2 = (f"{best['occupation']} achieves the best cost efficiency at £{best['cost_per_apply']:,.2f} per apply, "
                       f"while {worst['occupation']} is the most expensive at £{worst['cost_per_apply']:,.2f}.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2)}

    elif section == 'media':
        cat_stats = data.get('cat_stats')
        client_name = data.get('client_name', 'This client')

        if cat_stats is None or len(cat_stats) == 0:
            return {'intro': 'Media source data is not yet available for this client.',
                    'point_1': '', 'point_2': '', 'point_3': ''}

        sorted_stats = cat_stats.sort_values('total_applies', ascending=False)
        top_source = sorted_stats.iloc[0]
        intro = (f"{client_name}'s vacancies received traffic from {len(cat_stats)} distinct channels. "
                 f"{top_source['source_category']} is the leading source, generating "
                 f"{int(top_source['total_applies']):,} applies from {int(top_source['total_clicks']):,} views.")

        best_conv = cat_stats.loc[cat_stats['conversion_rate'].idxmax()]
        point_1 = ''
        if best_conv['conversion_rate'] > 0:
            point_1 = (f"{best_conv['source_category']} achieves the highest view-to-apply conversion rate at "
                       f"{best_conv['conversion_rate']:.1f}%, indicating well-matched candidates from this channel.")

        paid_rows = cat_stats[cat_stats['source_category'].str.contains('Paid|PPC|Sponsored', case=False, na=False)]
        point_2 = ''
        if len(paid_rows) > 0:
            paid = paid_rows.iloc[0]
            point_2 = (f"Paid channels ({paid['source_category']}) contributed {int(paid['total_applies']):,} applies "
                       f"with a {paid['conversion_rate']:.1f}% conversion rate.")

        point_3 = ''
        if len(cat_stats) > 1:
            second_source = sorted_stats.iloc[1] if len(sorted_stats) > 1 else None
            if second_source is not None:
                point_3 = (f"{second_source['source_category']} is the second strongest channel with "
                           f"{int(second_source['total_applies']):,} applies — providing diversified candidate flow.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    elif section == 'salary':
        per_occ = data.get('per_occ') or []
        client_name = data.get('client_name', 'This client')
        client_region = data.get('client_region')

        if not per_occ:
            return {'intro': 'Insufficient salary data to generate commentary for this client.',
                    'point_1': '', 'point_2': '', 'point_3': ''}

        # Compare each occupation's client mean to the national mean.
        # Tuples: (occupation, pct_diff_signed, client_mean, national_mean)
        deltas = []
        for p in per_occ:
            c, n = p.get('client_mean'), p.get('national_mean')
            if c is None or n is None or pd.isna(c) or pd.isna(n) or n == 0:
                continue
            deltas.append((p['occupation'], (c - n) / n * 100, c, n))

        above = sorted([d for d in deltas if d[1] > 0], key=lambda x: x[1], reverse=True)
        below = sorted([d for d in deltas if d[1] < 0], key=lambda x: x[1])  # most-negative first

        n_total = len(deltas)
        if n_total == 0:
            return {'intro': 'Salary data was insufficient to compute market comparisons.',
                    'point_1': '', 'point_2': '', 'point_3': ''}

        intro = (f"Across {client_name}'s top {n_total} most-posted occupations with salary data, "
                 f"{len(above)} pay above the national market average and {len(below)} pay below it.")

        point_1 = ''
        if above:
            top = above[0]
            point_1 = (f"{top[0]} is your strongest premium — sitting {top[1]:.0f}% above the national "
                       f"average (£{top[2]:,.0f} vs £{top[3]:,.0f}). Useful signal for attraction in this discipline.")

        point_2 = ''
        if below:
            worst = below[0]
            point_2 = (f"{worst[0]} sits {abs(worst[1]):.0f}% below the national average "
                       f"(£{worst[2]:,.0f} vs £{worst[3]:,.0f}) — a likely contributor to slower candidate flow in this category.")

        point_3 = ''
        if client_region:
            reg_above = reg_below = 0
            for p in per_occ:
                c, r = p.get('client_mean'), p.get('regional_mean')
                if c is None or r is None or pd.isna(c) or pd.isna(r):
                    continue
                if c > r:
                    reg_above += 1
                elif c < r:
                    reg_below += 1
            if reg_above + reg_below > 0:
                point_3 = (f"Within {client_region}, {client_name} pays above the regional average for "
                           f"{reg_above} of these roles and below for {reg_below} — useful context when "
                           f"benchmarking against employers competing for the same local talent pool.")
        else:
            point_3 = ("Regional benchmark unavailable for this client — common for central-government and "
                       "multi-site bodies whose vacancies span the UK.")

        return {'intro': _clean(intro), 'point_1': _clean(point_1), 'point_2': _clean(point_2), 'point_3': _clean(point_3)}

    return {'intro': 'No commentary available.', 'point_1': '', 'point_2': '', 'point_3': ''}


def render_client_report(df, media_df=None):
    """Render the Client Report tab — branded PPTX advertising report for renewals."""

    # PPTX placeholder geometry — surfaced in each chart caption so the user can
    # reason about text sizing against the actual rendered slot. Cached for an
    # hour; only re-reads `Renewals.pptx` if it changes.
    slot_dims = get_chart_slot_dimensions()

    # Determine the org column up-front; both the form and the renderer need it.
    org_col = 'organization_name' if 'organization_name' in df.columns else 'importer_name'

    # --- Settings state machine: form (open) <-> summary bar (collapsed) ---
    state = st.session_state
    state.setdefault('report_settings_collapsed', False)
    state.setdefault('report_generated', False)

    if not state['report_settings_collapsed']:
        # Settings form — all fields visible inside one bordered card, two per row
        with st.container(border=True):
            st.markdown('<h2 class="client-h2">Client advertising report</h2>', unsafe_allow_html=True)
            st.caption(
                "Pick a client and reporting period, then generate a branded report you can share or export."
            )

            col_client, col_dates = st.columns(2)
            with col_client:
                orgs = sorted(df[org_col].dropna().unique())
                orgs = [o for o in orgs if o and str(o).strip() not in ('', 'Unknown', 'nan')]
                org_counts = df.groupby(org_col).size().to_dict()
                st.selectbox(
                    "Select client",
                    orgs,
                    format_func=lambda name: f"{name} ({org_counts.get(name, 0):,} vacancies)",
                    key='report_client_name',
                )
            with col_dates:
                min_date = df['first_event_date'].dropna().min()
                max_date = df['last_event_date'].dropna().max()
                if pd.notna(min_date) and pd.notna(max_date):
                    min_d = min_date.date() if hasattr(min_date, 'date') else min_date
                    max_d = max_date.date() if hasattr(max_date, 'date') else max_date
                else:
                    min_d = datetime.now().date() - timedelta(days=365)
                    max_d = datetime.now().date()
                st.date_input("Report period", [min_d, max_d], key='report_dates')

            st.markdown('<div class="client-form-eyebrow">Cost &amp; report settings</div>', unsafe_allow_html=True)
            cost_col1, cost_col2 = st.columns(2)
            with cost_col1:
                st.number_input(
                    "Annual spend (GBP)", value=0.0, step=100.0, format="%.2f",
                    key='report_spend',
                    help="Enter 0 to skip the ROI section",
                )
            with cost_col2:
                st.number_input(
                    "Rate card price per job (GBP)", value=600.0, step=10.0, format="%.2f",
                    key='report_rate_card',
                )
            st.checkbox(
                "Include selected client in benchmark",
                value=False,
                key='report_include_self',
                help="When unchecked, the selected client is excluded from the benchmark average (recommended for fair comparison)",
            )

            st.markdown('<div class="client-form-eyebrow">Contact details <span class="client-form-eyebrow-meta">(for PDF)</span></div>', unsafe_allow_html=True)
            contact_col1, contact_col2 = st.columns(2)
            with contact_col1:
                st.text_input("Account manager", key='report_contact_name', placeholder="e.g. Jane Smith")
            with contact_col2:
                st.text_input("Title", key='report_contact_title', placeholder="e.g. Account Director")
            contact_col3, contact_col4 = st.columns(2)
            with contact_col3:
                st.text_input("Email", key='report_contact_email', placeholder="e.g. jane@jgp.co.uk")
            with contact_col4:
                st.text_input("Phone", key='report_contact_phone', placeholder="e.g. 020 7946 0958")

            if st.button("Generate report", type="primary", key='report_generate'):
                # Pin widget values to non-widget keys. Streamlit Cloud
                # garbage-collects state entries for widgets that aren't
                # rendered on a given run — so once the settings card
                # collapses, the next rerun (e.g. when Download is clicked)
                # can clear state['report_client_name'] and friends.
                # Reading from _applied_* keys keeps the report stable.
                state['_applied_client'] = state.get('report_client_name')
                state['_applied_dates'] = state.get('report_dates')
                state['_applied_spend'] = state.get('report_spend', 0.0)
                state['_applied_rate_card'] = state.get('report_rate_card', 600.0)
                state['_applied_include_self'] = bool(state.get('report_include_self', False))
                state['_applied_contact_name'] = state.get('report_contact_name', '') or ''
                state['_applied_contact_title'] = state.get('report_contact_title', '') or ''
                state['_applied_contact_email'] = state.get('report_contact_email', '') or ''
                state['_applied_contact_phone'] = state.get('report_contact_phone', '') or ''
                state['report_settings_collapsed'] = True
                state['report_generated'] = True
                st.rerun()

    if not state['report_generated']:
        st.info("Pick a client and click **Generate report** to build the advertising report.")
        return

    # --- Read pinned form values ---
    selected_client = state.get('_applied_client')
    if not selected_client:
        st.warning("No client selected. Open settings and pick one.")
        return
    report_dates = state.get('_applied_dates', [])
    annual_spend = float(state.get('_applied_spend', 0.0) or 0.0)
    rate_card_price = float(state.get('_applied_rate_card', 600.0) or 0.0)
    include_self = bool(state.get('_applied_include_self', False))
    contact_name = state.get('_applied_contact_name', '') or ''
    contact_title = state.get('_applied_contact_title', '') or ''
    contact_email = state.get('_applied_contact_email', '') or ''
    contact_phone = state.get('_applied_contact_phone', '') or ''

    if not report_dates or len(report_dates) < 2:
        st.warning("Please select a start and end date in settings.")
        return

    report_start, report_end = report_dates[0], report_dates[1]
    period_str = f"{report_start.strftime('%d %b %Y')} – {report_end.strftime('%d %b %Y')}"

    # --- Summary bar (collapsed-settings state) ---
    if state['report_settings_collapsed']:
        col_summary, col_edit = st.columns([6, 1])
        with col_summary:
            st.markdown(
                summary_bar(selected_client, period_str, am_name=contact_name),
                unsafe_allow_html=True,
            )
        with col_edit:
            if st.button("Edit", key='report_edit'):
                # Re-seed widget keys from the pinned values so the form
                # opens pre-filled (Streamlit reads widget initial value
                # from session_state at the start of the run that renders
                # the widget).
                state['report_client_name'] = state.get('_applied_client')
                state['report_dates'] = state.get('_applied_dates')
                state['report_spend'] = state.get('_applied_spend', 0.0)
                state['report_rate_card'] = state.get('_applied_rate_card', 600.0)
                state['report_include_self'] = bool(state.get('_applied_include_self', False))
                state['report_contact_name'] = state.get('_applied_contact_name', '')
                state['report_contact_title'] = state.get('_applied_contact_title', '')
                state['report_contact_email'] = state.get('_applied_contact_email', '')
                state['report_contact_phone'] = state.get('_applied_contact_phone', '')
                state['report_settings_collapsed'] = False
                st.rerun()

    # --- Data preparation ---

    # Client data — filter on organization_name (or importer_name fallback)
    client_df = df[df[org_col] == selected_client].copy()
    client_df = client_df[
        (client_df['last_event_date'].dt.date >= report_start) &
        (client_df['first_event_date'].dt.date <= report_end)
    ]

    if len(client_df) == 0:
        st.warning(f"No vacancies found for **{selected_client}** in the selected date range.")
        return

    # Benchmark = market data in same date range (exclude self by default for fair comparison)
    bench_mask = (
        (df['last_event_date'].dt.date >= report_start) &
        (df['first_event_date'].dt.date <= report_end)
    )
    if not include_self:
        bench_mask = bench_mask & (df[org_col] != selected_client)
    benchmark_df = df[bench_mask].copy()

    # Media data for client — match via entity_id (most reliable link)
    # Note: prepare_enriched_data() renames entity_id_str → entity_id in df,
    # but media_df keeps the original entity_id_str column name.
    client_media = None
    if media_df is not None and len(media_df) > 0:
        client_eid_col = 'entity_id' if 'entity_id' in client_df.columns else 'entity_id_str'
        media_eid_col = 'entity_id_str' if 'entity_id_str' in media_df.columns else 'entity_id'
        if client_eid_col in client_df.columns and media_eid_col in media_df.columns:
            client_entity_ids = client_df[client_eid_col].dropna().unique()
            client_media = media_df[media_df[media_eid_col].isin(client_entity_ids)].copy()
        if client_media is None or len(client_media) == 0:
            client_media = None
        else:
            client_media = apply_media_categories(client_media)

    # --- Hero band ---
    # Dynamic lede: vacancies, views vs benchmark, applies vs benchmark,
    # blended CPA, and rate-card saving. Each clause is dropped if the
    # underlying data is missing (no benchmark, no spend, no saving).
    hero_num_jobs = len(client_df)
    hero_total_applies = int(client_df['applies'].sum())
    hero_rate_card_total = rate_card_price * hero_num_jobs
    hero_avg_clicks = client_df['clicks'].mean() if hero_num_jobs > 0 else 0
    hero_avg_applies = client_df['applies'].mean() if hero_num_jobs > 0 else 0

    if len(benchmark_df) > 0:
        bench_avg_clicks = benchmark_df['clicks'].mean()
        bench_avg_applies = benchmark_df['applies'].mean()
        hero_views_pct = (hero_avg_clicks / bench_avg_clicks * 100) if bench_avg_clicks > 0 else None
        hero_applies_pct = (hero_avg_applies / bench_avg_applies * 100) if bench_avg_applies > 0 else None
    else:
        bench_avg_clicks = 0
        bench_avg_applies = 0
        hero_views_pct = None
        hero_applies_pct = None

    lede_intro = (
        f"Across <strong>{hero_num_jobs:,} vacancies</strong> in this period, "
        f"{selected_client}"
    )
    lede_perf = ""
    if hero_views_pct is not None and hero_applies_pct is not None:
        views_delta = hero_views_pct - 100
        applies_delta = hero_applies_pct - 100
        if views_delta >= 0 and applies_delta >= 0:
            verb = "outperformed the market"
        elif views_delta < 0 and applies_delta < 0:
            verb = "underperformed the market"
        else:
            verb = "saw mixed market performance"
        lede_perf = (
            f" {verb} on visibility "
            f"(<strong>{views_delta:+.0f}% views</strong>) "
            f"and engagement "
            f"(<strong>{applies_delta:+.0f}% applies</strong>)"
        )

    lede_cost = ""
    if annual_spend > 0 and hero_total_applies > 0:
        cpa = annual_spend / hero_total_applies
        lede_cost = f", with a blended cost per apply of <strong>£{cpa:,.2f}</strong>"
        if hero_rate_card_total > 0:
            saving = hero_rate_card_total - annual_spend
            if saving > 0:
                saving_pct = saving / hero_rate_card_total * 100
                lede_cost += (
                    f" and a <strong>{saving_pct:.0f}% saving</strong> "
                    "vs rate-card pricing"
                )

    lede = (lede_intro + lede_perf + lede_cost + ".").strip()

    st.markdown(
        client_hero(selected_client, period_str, lede, hero_num_jobs, am_name=contact_name),
        unsafe_allow_html=True,
    )

    # Store all figures for PDF export
    report_figures = {}

    # ===================================================================
    # SECTION 01: HEADLINE NUMBERS
    # ===================================================================

    # KPI 4 (Cost per apply) helper falls back gracefully when spend is 0
    # or there are no applies yet — the value itself goes to em-dash.
    if annual_spend > 0 and hero_total_applies > 0:
        headline_cpa = f"£{annual_spend / hero_total_applies:,.2f}"
        headline_cpa_helper = f"£{annual_spend:,.0f} total spend"
    elif annual_spend > 0:
        headline_cpa = "—"
        headline_cpa_helper = f"£{annual_spend:,.0f} total spend, no applies yet"
    else:
        headline_cpa = "—"
        headline_cpa_helper = "Set annual spend in settings"

    # KPI 2 + 3 helpers carry the benchmark delta + raw benchmark mean
    if hero_views_pct is not None:
        views_delta = hero_views_pct - 100
        helper_views = (
            f"{views_delta:+.0f}% vs benchmark "
            f"({bench_avg_clicks:,.0f} views)"
        )
    else:
        helper_views = "No market benchmark available"

    if hero_applies_pct is not None:
        applies_delta = hero_applies_pct - 100
        helper_applies = (
            f"{applies_delta:+.0f}% vs benchmark "
            f"({bench_avg_applies:,.1f} applies)"
        )
    else:
        helper_applies = "No market benchmark available"

    with st.container(border=True):
        st.markdown(
            section_eyebrow('01', 'Headline numbers', short='Headlines')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_01_headlines"]}</p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )
        st.markdown(
            '<div class="client-kpi-grid">'
            + kpi_card_dark(
                "Jobs advertised", f"{hero_num_jobs:,}",
                helper="Across the reporting period",
            )
            + kpi_card(
                "Avg views per vacancy", f"{hero_avg_clicks:,.0f}",
                helper=helper_views,
            )
            + kpi_card(
                "Avg applies per vacancy", f"{hero_avg_applies:,.1f}",
                helper=helper_applies,
            )
            + kpi_card("Cost per apply", headline_cpa, helper=headline_cpa_helper)
            + '</div>',
            unsafe_allow_html=True,
        )

    # ===================================================================
    # SECTION 02: PER-VACANCY BENCHMARKING (existing scatter)
    # ===================================================================
    with st.container(border=True):
        st.markdown(
            section_eyebrow('02', 'Per-vacancy benchmarking', short='Benchmarking')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_02_intro"]}'
            + f' <span class="info-icon" data-tooltip="{CHART_EXPLAINERS["benchmark_scatter"]}">?</span></p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )

        # Calculate per-occupation benchmark averages (from ALL clients)
        occ_benchmarks = benchmark_df.groupby('occupation').agg(
            avg_clicks=('clicks', 'mean'),
            avg_applies=('applies', 'mean'),
            vacancy_count=('clicks', 'count')
        ).reset_index()

        # Only use occupations with enough data for reliable benchmarks
        MIN_BENCHMARK_VACANCIES = 5
        reliable_occs = occ_benchmarks[occ_benchmarks['vacancy_count'] >= MIN_BENCHMARK_VACANCIES]

        # Calculate % difference from benchmark for each client vacancy (vectorized)
        scatter_df = client_df[['title', 'occupation', 'clicks', 'applies']].copy()
        scatter_df = scatter_df.merge(
            reliable_occs[['occupation', 'avg_clicks', 'avg_applies']],
            on='occupation', how='left'
        )

        # Vectorized % diff calculations (safe division)
        scatter_df['views_diff_pct'] = (
            (scatter_df['clicks'] - scatter_df['avg_clicks'])
            / scatter_df['avg_clicks'].replace(0, np.nan) * 100
        ).fillna(0)
        scatter_df['applies_diff_pct'] = (
            (scatter_df['applies'] - scatter_df['avg_applies'])
            / scatter_df['avg_applies'].replace(0, np.nan) * 100
        ).fillna(0)

        # Categorize vacancies into 4 groups (matching original PDF categories)
        has_benchmark = scatter_df['avg_clicks'].notna()
        zero_applies = scatter_df['applies'] == 0
        occ_median_clicks = scatter_df.loc[has_benchmark, 'clicks'].median() if has_benchmark.any() else 0
        high_traffic = scatter_df['clicks'] > occ_median_clicks

        scatter_df['category'] = np.select(
            [has_benchmark & ~zero_applies,
             has_benchmark & zero_applies & high_traffic,
             has_benchmark & zero_applies & ~high_traffic],
            ['Benchmarkable', 'Zero Applies (Possible Redirect)', 'Zero Applies (Low Traffic)'],
            default='Low Sample (No Benchmark)'
        )

        # Clean up merge columns and fill NaN diffs for no-benchmark rows
        scatter_df.drop(columns=['avg_clicks', 'avg_applies'], inplace=True)
        scatter_df.loc[~has_benchmark, ['views_diff_pct', 'applies_diff_pct']] = 0

        # --- Brand colour + marker maps for 4 categories ---
        category_colors = {
            'Benchmarkable': JGP_COLORS['primary'],
            'Zero Applies (Possible Redirect)': JGP_COLORS['blue'],
            'Zero Applies (Low Traffic)': JGP_COLORS['negative'],
            'Low Sample (No Benchmark)': JGP_COLORS['light_purple'],
        }
        category_symbols = {
            'Benchmarkable': 'triangle-up',
            'Zero Applies (Possible Redirect)': 'diamond',
            'Zero Applies (Low Traffic)': 'x',
            'Low Sample (No Benchmark)': 'circle',
        }

        chart_col, grid_col = st.columns([3, 2])

        with chart_col:
            if len(scatter_df) > 0:
                fig_scatter = go.Figure()
                for cat in ['Benchmarkable', 'Zero Applies (Possible Redirect)', 'Zero Applies (Low Traffic)', 'Low Sample (No Benchmark)']:
                    cat_data = scatter_df[scatter_df['category'] == cat]
                    if len(cat_data) == 0:
                        continue
                    fig_scatter.add_trace(go.Scatter(
                        x=cat_data['applies_diff_pct'],
                        y=cat_data['views_diff_pct'],
                        mode='markers',
                        name=cat,
                        marker=dict(
                            color=category_colors[cat],
                            symbol=category_symbols[cat],
                            size=10,
                            line=dict(width=1, color='white')
                        ),
                        customdata=cat_data[['title', 'occupation', 'clicks', 'applies']].values,
                        hovertemplate=(
                            "<b>%{customdata[0]}</b><br>"
                            "Occupation: %{customdata[1]}<br>"
                            "Views: %{customdata[2]}<br>"
                            "Applies: %{customdata[3]}<br>"
                            "Views vs Bench: %{y:+.0f}%<br>"
                            "Applies vs Bench: %{x:+.0f}%<extra></extra>"
                        ),
                    ))
                fig_scatter.add_hline(y=0, line_dash="dash", line_color="grey", opacity=0.5)
                fig_scatter.add_vline(x=0, line_dash="dash", line_color="grey", opacity=0.5)
                fig_scatter.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
                fig_scatter.update_layout(
                    height=500,
                    showlegend=True,
                    legend=dict(orientation='h', y=-0.2, font=dict(size=11)),
                    margin=dict(t=20, b=40, l=40, r=20),
                    xaxis_title="Applies Difference from Benchmark (%)",
                    yaxis_title="Views Difference from Benchmark (%)",
                    plot_bgcolor='rgba(0,0,0,0)',
                    xaxis=dict(gridcolor=JGP_COLORS['light_purple'], gridwidth=1, zeroline=False),
                    yaxis=dict(gridcolor=JGP_COLORS['light_purple'], gridwidth=1, zeroline=False),
                )
                st.plotly_chart(fig_scatter, use_container_width=True, config={'displayModeBar': 'hover', 'displaylogo': False})
                report_figures['scatter'] = fig_scatter

        benchmarkable_count = len(scatter_df[scatter_df['category'] == 'Benchmarkable'])
        zero_redirect = len(scatter_df[scatter_df['category'] == 'Zero Applies (Possible Redirect)'])
        zero_low = len(scatter_df[scatter_df['category'] == 'Zero Applies (Low Traffic)'])
        no_bench_count = len(scatter_df[scatter_df['category'] == 'Low Sample (No Benchmark)'])

        # Pre-compute commentary so we can render it inside grid_col under the
        # status grid (rather than full-width below the chart + grid row).
        benchmarkable_rows = scatter_df[scatter_df['category'] == 'Benchmarkable'].copy()
        top_performers = []
        worst_performers = []
        if len(benchmarkable_rows) > 0:
            benchmarkable_rows['_score'] = benchmarkable_rows['views_diff_pct'] + benchmarkable_rows['applies_diff_pct']
            top_sorted = benchmarkable_rows.nlargest(3, '_score')
            top_performers = top_sorted[['title', 'occupation', 'views_diff_pct', 'applies_diff_pct']].to_dict('records')
            worst_sorted = benchmarkable_rows.nsmallest(3, '_score')
            worst_performers = worst_sorted[['title', 'occupation', 'views_diff_pct', 'applies_diff_pct']].to_dict('records')

        commentary = generate_section_commentary('scatter', {
            'total_count': len(scatter_df),
            'benchmarkable_count': benchmarkable_count,
            'zero_applies_count': zero_redirect + zero_low,
            'no_benchmark_count': no_bench_count,
            'top_performers': top_performers,
            'worst_performers': worst_performers,
        })

        with grid_col:
            st.markdown(
                status_grid([
                    {'label': 'Benchmarkable',     'value': f"{benchmarkable_count:,}",
                     'help': 'In-period vs occupation average'},
                    {'label': 'Low traffic',       'value': f"{zero_low:,}",
                     'help': 'Zero applies, below median views'},
                    {'label': 'Possible redirect', 'value': f"{zero_redirect:,}",
                     'help': 'Zero applies, high views'},
                    {'label': 'No benchmark',      'value': f"{no_bench_count:,}",
                     'help': 'Occupation lacks ≥5 market peers'},
                ]),
                unsafe_allow_html=True,
            )
            st.markdown(commentary_panel(commentary), unsafe_allow_html=True)

    # ===================================================================
    # SECTION 03: PERFORMANCE VS MARKET (existing benchmarking summary)
    # ===================================================================
    with st.container(border=True):
        st.markdown(
            section_eyebrow('03', 'Performance vs market', short='Performance')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_03_intro"]}'
            + f' <span class="info-icon" data-tooltip="{CHART_EXPLAINERS["benchmark_average"]}">?</span></p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )

        benchmark_avg_clicks = benchmark_df['clicks'].mean() if len(benchmark_df) > 0 else 0
        benchmark_avg_applies = benchmark_df['applies'].mean() if len(benchmark_df) > 0 else 0
        client_avg_clicks = client_df['clicks'].mean() if len(client_df) > 0 else 0
        client_avg_applies = client_df['applies'].mean() if len(client_df) > 0 else 0
        views_pct = (client_avg_clicks / benchmark_avg_clicks * 100) if benchmark_avg_clicks > 0 else 0
        applies_pct = (client_avg_applies / benchmark_avg_applies * 100) if benchmark_avg_applies > 0 else 0

        bench_commentary = generate_section_commentary('benchmark', {
            'client_avg_clicks': client_avg_clicks,
            'benchmark_avg_clicks': benchmark_avg_clicks,
            'client_avg_applies': client_avg_applies,
            'benchmark_avg_applies': benchmark_avg_applies,
            'num_jobs': len(client_df),
            'client_name': selected_client,
        })

        bench_chart_col, bench_stats_col = st.columns([3, 2])

        with bench_chart_col:
            # Indexed bars at 100% baseline
            fig_bench = go.Figure()
            fig_bench.add_trace(go.Bar(
                x=['Views', 'Applies'],
                y=[views_pct, applies_pct],
                marker_color=[
                    JGP_COLORS['positive'] if views_pct >= 100 else JGP_COLORS['blue'],
                    JGP_COLORS['positive'] if applies_pct >= 100 else JGP_COLORS['blue'],
                ],
                text=[f"{views_pct:.0f}%", f"{applies_pct:.0f}%"],
                textposition='outside',
                textfont=dict(size=14, color=JGP_COLORS['deep_blue']),
            ))
            fig_bench.add_hline(
                y=100, line_dash="dash", line_width=3, line_color=JGP_COLORS['primary'],
                annotation_text="Market benchmark (100%)",
                annotation_position="top right",
                annotation_bgcolor="white",
                annotation_bordercolor=JGP_COLORS['deep_blue'],
                annotation_borderwidth=1,
                annotation_font=dict(color=JGP_COLORS['deep_blue'], size=13),
            )
            fig_bench.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_bench.update_layout(
                height=500,
                showlegend=False,
                margin=dict(t=20, b=40, l=40, r=20),
                yaxis_title="% of market benchmark",
                yaxis_range=[0, max(views_pct, applies_pct, 100) * 1.25],
            )
            st.plotly_chart(
                fig_bench,
                use_container_width=True,
                config={'displayModeBar': 'hover', 'displaylogo': False},
            )
            report_figures['benchmark_combined'] = fig_bench

        with bench_stats_col:
            st.markdown(
                status_grid(
                    [
                        {
                            'label': 'Average views per vacancy',
                            'value': f"{client_avg_clicks:,.0f}",
                            'help': f"Market benchmark {benchmark_avg_clicks:,.0f} ({views_pct:.0f}% of market)",
                        },
                        {
                            'label': 'Average applies per vacancy',
                            'value': f"{client_avg_applies:,.1f}",
                            'help': f"Market benchmark {benchmark_avg_applies:,.1f} ({applies_pct:.0f}% of market)",
                        },
                    ],
                ),
                unsafe_allow_html=True,
            )
            st.markdown(commentary_panel(bench_commentary), unsafe_allow_html=True)

    # ===================================================================
    # SECTION 04: POSTINGS & APPLY VOLUME (existing job postings)
    # ===================================================================
    with st.container(border=True):
        st.markdown(
            section_eyebrow('04', 'Postings & apply volume', short='Postings')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_04_intro"]}'
            + f' <span class="info-icon" data-tooltip="{CHART_EXPLAINERS["postings_by_type"]}">?</span></p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )

        # Client view per occupation.
        by_type = client_df.groupby('occupation').agg(
            jobs_posted=('clicks', 'count'),
            apply_clicks=('applies', 'sum'),
        ).reset_index()
        by_type = by_type[by_type['jobs_posted'] >= 1]
        by_type['applies_per_posting'] = by_type['apply_clicks'] / by_type['jobs_posted']

        # Market benchmark applies-per-posting per occupation (mean across all
        # other clients' postings in the same occupation).
        bench_by_occ = benchmark_df.groupby('occupation').agg(
            bench_jobs=('clicks', 'count'),
            bench_applies=('applies', 'sum'),
        ).reset_index()
        bench_by_occ['bench_applies_per_post'] = (
            bench_by_occ['bench_applies'] / bench_by_occ['bench_jobs']
        )
        by_type = by_type.merge(
            bench_by_occ[['occupation', 'bench_applies_per_post']],
            on='occupation',
            how='left',
        )

        # Table sort order: most-posted occupations first (descending). The
        # chart still uses ascending so the horizontal bars read bottom-up.
        by_type_sorted = by_type.sort_values(
            'jobs_posted', ascending=False
        ).reset_index(drop=True)

        by_type_top10 = by_type_sorted.head(10)
        items_html = ''
        for row in by_type_top10.itertuples():
            bench_str = (
                '—' if pd.isna(row.bench_applies_per_post)
                else f"{row.bench_applies_per_post:.2f}"
            )
            items_html += (
                '<li>'
                f'<span>{row.occupation}</span>'
                f'<span class="num">{int(row.jobs_posted):,}</span>'
                f'<span class="num">{int(row.apply_clicks):,}</span>'
                f'<span class="num">{row.applies_per_posting:.2f}</span>'
                f'<span class="num">{bench_str}</span>'
                '</li>'
            )
        header_html = (
            '<div class="rank-list__header">'
            '<span></span>'
            '<span>Occupation</span>'
            '<span class="num">Postings</span>'
            '<span class="num">Applies</span>'
            '<span class="num">A/posting</span>'
            '<span class="num">Bench A/p</span>'
            '</div>'
        )
        table_html = (
            '<div class="rank-list">'
            '<h4>Top Occupations for Applies vs posts</h4>'
            f'{header_html}'
            f'<ol class="rank-list__items">{items_html}</ol>'
            '</div>'
        )

        total_jobs = len(client_df)
        total_applies_val = int(client_df['applies'].sum())
        postings_commentary = generate_section_commentary('postings', {
            'total_jobs': total_jobs,
            'total_applies': total_applies_val,
            'by_type': by_type_sorted,
            'client_name': selected_client,
        })

        postings_table_col, postings_comm_col = st.columns([3, 2])
        with postings_table_col:
            st.markdown(table_html, unsafe_allow_html=True)
        with postings_comm_col:
            st.markdown(commentary_panel(postings_commentary), unsafe_allow_html=True)

        with st.expander("Full chart breakdown", expanded=False):
            # Chart still sorts ascending so horizontal bars build bottom-up.
            by_type_chart = by_type.sort_values('jobs_posted', ascending=True)
            fig_postings = go.Figure()
            fig_postings.add_trace(go.Bar(
                y=by_type_chart['occupation'], x=by_type_chart['jobs_posted'],
                name='Postings', orientation='h',
                marker_color=JGP_COLORS['primary'],
                text=by_type_chart['jobs_posted'], textposition='outside',
            ))
            fig_postings.add_trace(go.Bar(
                y=by_type_chart['occupation'], x=by_type_chart['apply_clicks'],
                name='Applies', orientation='h',
                marker_color=JGP_COLORS['accent'],
                text=by_type_chart['apply_clicks'].astype(int), textposition='outside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_postings.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_postings.update_layout(
                height=max(500, len(by_type_chart) * 40),
                margin=dict(t=20, b=40, l=40, r=20),
                barmode='group',
                legend=dict(orientation='h', y=-0.1),
                xaxis_title="Count", yaxis_title="",
                bargap=0.1, bargroupgap=0.0,
            )
            st.plotly_chart(
                fig_postings,
                use_container_width=True,
                config={'displayModeBar': 'hover', 'displaylogo': False},
            )
            report_figures['postings'] = fig_postings

    # ===================================================================
    # SECTION 05: ADVERTISING ROI (existing; CPA split out in section 06 later)
    # ===================================================================
    with st.container(border=True):
        st.markdown(
            section_eyebrow('05', 'Advertising ROI', short='ROI')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_05_intro"]}'
            + f' <span class="info-icon" data-tooltip="{CHART_EXPLAINERS["spend_vs_ratecard"]}">?</span></p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )

        num_jobs = len(client_df)
        total_clicks = int(client_df['clicks'].sum())
        total_applies_val = int(client_df['applies'].sum())

        roi_by_type = None
        if annual_spend > 0:
            cost_per_job = annual_spend / num_jobs if num_jobs > 0 else 0
            cost_per_view = annual_spend / total_clicks if total_clicks > 0 else 0
            cost_per_apply = annual_spend / total_applies_val if total_applies_val > 0 else 0
            rate_card_total = rate_card_price * num_jobs
            saving_pct = ((rate_card_total - annual_spend) / rate_card_total * 100) if rate_card_total > 0 else 0
            saving_amount = max(0, rate_card_total - annual_spend)

            # Four KPIs — one filled deep-blue (the headline saving)
            roi_kpi1, roi_kpi2, roi_kpi3, roi_kpi4 = st.columns(4)
            with roi_kpi1:
                st.markdown(
                    kpi_card("Annual spend", f"£{annual_spend:,.0f}",
                             helper="As entered in settings"),
                    unsafe_allow_html=True,
                )
            with roi_kpi2:
                st.markdown(
                    kpi_card("Rate card equivalent", f"£{rate_card_total:,.0f}",
                             helper=f"£{rate_card_price:,.0f} list × {num_jobs:,} ads"),
                    unsafe_allow_html=True,
                )
            with roi_kpi3:
                st.markdown(
                    kpi_card_dark("Saving vs rate card", f"£{saving_amount:,.0f}",
                                  helper=f"{saving_pct:.0f}% under list"),
                    unsafe_allow_html=True,
                )
            with roi_kpi4:
                st.markdown(
                    kpi_card("Cost per view", f"£{cost_per_view:,.2f}",
                             helper="Annual spend ÷ total views"),
                    unsafe_allow_html=True,
                )

            # Saving bar — your spend stacked with the saving (green) makes
            # the rate-card equivalent the full bar.
            fig_roi = go.Figure()
            fig_roi.add_trace(go.Bar(
                y=['Cost'], x=[annual_spend],
                name='Your spend',
                orientation='h',
                marker_color=JGP_COLORS['primary'],
                text=[f"£{annual_spend:,.0f}"], textposition='inside',
                textfont=dict(color=JGP_COLORS['white']),
            ))
            fig_roi.add_trace(go.Bar(
                y=['Cost'], x=[saving_amount],
                name=f'Saving ({saving_pct:.0f}%)',
                orientation='h',
                marker_color=JGP_COLORS['accent'],
                text=[f"£{saving_amount:,.0f}"], textposition='inside',
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_roi.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_roi.update_layout(
                barmode='stack',
                height=180,
                showlegend=True,
                legend=dict(orientation='h', y=-0.4),
                xaxis_title="GBP",
                yaxis_title="",
                margin=dict(t=20, b=40, l=40, r=20),
            )
            st.plotly_chart(
                fig_roi,
                use_container_width=True,
                config={'displayModeBar': 'hover', 'displaylogo': False},
            )
            report_figures['roi_cost'] = fig_roi

            # Cost-per-apply by type — used in section 06 below + commentary
            roi_by_type = client_df.groupby('occupation').agg(
                total_applies=('applies', 'sum'),
                job_count=('clicks', 'count')
            ).reset_index()
            roi_by_type = roi_by_type[roi_by_type['total_applies'] > 0]
            roi_by_type['cost_allocated'] = annual_spend * (roi_by_type['job_count'] / roi_by_type['job_count'].sum())
            roi_by_type['cost_per_apply'] = roi_by_type['cost_allocated'] / roi_by_type['total_applies']
            roi_by_type = roi_by_type.sort_values('cost_per_apply', ascending=True)

            roi_commentary = generate_section_commentary('roi', {
                'annual_spend': annual_spend,
                'rate_card_price': rate_card_price,
                'num_jobs': num_jobs,
                'total_clicks': total_clicks,
                'total_applies': total_applies_val,
                'cost_per_job': cost_per_job,
                'cost_per_view': cost_per_view,
                'cost_per_apply': cost_per_apply,
                'saving_pct': saving_pct,
                'roi_by_type': roi_by_type if len(roi_by_type) > 0 else None,
            })
            st.markdown(commentary_panel(roi_commentary), unsafe_allow_html=True)
        else:
            st.markdown(
                kpi_card_dark("Vacancies advertised", f"{num_jobs:,}",
                              helper="Set annual spend in settings to see ROI analysis"),
                unsafe_allow_html=True,
            )

    # ===================================================================
    # SECTION 06: COST PER APPLY (cheapest/priciest + collapsible breakdown)
    # ===================================================================
    if annual_spend > 0 and roi_by_type is not None and len(roi_by_type) > 0:
        with st.container(border=True):
            st.markdown(
                section_eyebrow('06', 'Cost per apply', short='Cost per apply')
                + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_06_intro"]}'
                + f' <span class="info-icon" data-tooltip="{CHART_EXPLAINERS["cost_per_app_by_occupation"]}">?</span></p>'
                + '<hr class="client-section-divider" />',
                unsafe_allow_html=True,
            )

            cheapest = roi_by_type.head(3)
            priciest = roi_by_type.tail(3).iloc[::-1]

            col_cheap, col_pricey = st.columns(2)
            with col_cheap:
                cheapest_html = '<div class="cpa-list"><h4>Cheapest occupations</h4><ol>'
                for _, row in cheapest.iterrows():
                    cheapest_html += (
                        f'<li><span class="cpa-label">{row["occupation"]}</span>'
                        f'<span class="cpa-value">£{row["cost_per_apply"]:,.2f}</span></li>'
                    )
                cheapest_html += '</ol></div>'
                st.markdown(cheapest_html, unsafe_allow_html=True)
            with col_pricey:
                pricey_html = '<div class="cpa-list"><h4>Priciest occupations</h4><ol>'
                for _, row in priciest.iterrows():
                    pricey_html += (
                        f'<li><span class="cpa-label">{row["occupation"]}</span>'
                        f'<span class="cpa-value">£{row["cost_per_apply"]:,.2f}</span></li>'
                    )
                pricey_html += '</ol></div>'
                st.markdown(pricey_html, unsafe_allow_html=True)

            with st.expander("Full cost-per-apply breakdown", expanded=False):
                fig_cpa = go.Figure()
                fig_cpa.add_trace(go.Bar(
                    y=roi_by_type['occupation'], x=roi_by_type['cost_per_apply'],
                    orientation='h', marker_color=JGP_COLORS['primary'],
                    text=roi_by_type['cost_per_apply'].apply(lambda x: f"£{x:,.2f}"),
                    textposition='outside'
                ))
                fig_cpa.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
                fig_cpa.update_layout(
                    height=max(300, len(roi_by_type) * 35),
                    margin=dict(t=20, b=40, l=40, r=20),
                    xaxis_title="Cost per apply (GBP)", yaxis_title="",
                )
                st.plotly_chart(
                    fig_cpa,
                    use_container_width=True,
                    config={'displayModeBar': 'hover', 'displaylogo': False},
                )
                report_figures['roi_cpa'] = fig_cpa

    # ===================================================================
    # SECTION 07: SALARY BENCHMARKS (existing top-10 occupations)
    # ===================================================================
    # For each of the client's most-posted-with-salary roles, show the
    # market salary distribution as a histogram and overlay three reference
    # means: client, national, regional (client's HQ region). Mirrors the
    # salary-tab histogram pattern (views/salary.py:185-239) per occupation.
    with st.container(border=True):
        st.markdown(
            section_eyebrow('07', 'Salary benchmarks', short='Salaries')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_07_intro"]}'
            + f' <span class="info-icon" data-tooltip="{CHART_EXPLAINERS["salary_by_occupation"]}">?</span></p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )

        # Persistent across the conditional branches so the commentary generator
        # downstream can pick them up (None when section is skipped).
        salary_per_occ = None
        salary_client_region = None

        client_with_salary = client_df[client_df.get('has_salary_data', False) == True]

        if len(client_with_salary) == 0:
            st.info(
                "No salary data is available for this client's vacancies in the "
                "selected period — salary benchmark omitted."
            )
        else:
            occ_counts = client_with_salary['occupation'].dropna().value_counts()
            qualifying = occ_counts[occ_counts >= 5]
            top_occupations = qualifying.head(10).index.tolist()

            if len(top_occupations) == 0:
                st.info(
                    "No occupations have at least 5 vacancies with salary data for "
                    "this client in the selected period — salary benchmark needs "
                    "≥5 priced roles per occupation to be meaningful."
                )
            else:
                # Look up client HQ region (None for multi-site / central-gov clients)
                hq_map = load_client_hq_regions()
                client_region = hq_map.get(selected_client.lower().strip())

                # Pre-build the regional market subset once. Compare on
                # lower-stripped strings so canonical/raw spelling differences
                # between primary_uk_region and client_hq_addresses don't drop
                # the line silently.
                df_regional_market = None
                if client_region and 'primary_uk_region' in df.columns:
                    norm = client_region.strip().lower()
                    df_regional_market = df[
                        (df.get('has_salary_data', False) == True)
                        & (df['primary_uk_region'].fillna('').str.strip().str.lower() == norm)
                    ]
                    if len(df_regional_market) == 0:
                        df_regional_market = None  # No samples → drop regional line

                # Brand-kit reference lines: pink (you), bold green (national), deep blue (region)
                client_color = JGP_COLORS['pink']        # #ffc4c4 — your mean
                national_color = JGP_COLORS['accent']    # #e5ff6e — national mean
                regional_color = JGP_COLORS['deep_blue'] # #240f45 — regional mean

                per_occ = []
                for occ in top_occupations:
                    client_occ = client_with_salary[client_with_salary['occupation'] == occ]
                    client_mean = client_occ['annual_mid_salary'].mean()

                    market_occ = df[(df['occupation'] == occ) & (df.get('has_salary_data', False) == True)]
                    market_salaries = market_occ['annual_mid_salary'].dropna()
                    national_mean = market_salaries.mean() if len(market_salaries) else np.nan

                    if df_regional_market is not None:
                        reg_vals = df_regional_market[df_regional_market['occupation'] == occ]['annual_mid_salary'].dropna()
                        regional_mean = reg_vals.mean() if len(reg_vals) >= 3 else np.nan
                        regional_n = len(reg_vals)
                    else:
                        regional_mean = np.nan
                        regional_n = 0

                    per_occ.append({
                        'occupation': occ,
                        'client_n': len(client_occ),
                        'market_n': len(market_salaries),
                        'regional_n': regional_n,
                        'market_salaries': market_salaries,
                        'client_mean': client_mean,
                        'national_mean': national_mean,
                        'regional_mean': regional_mean,
                    })

                any_regional = any(not pd.isna(p['regional_mean']) for p in per_occ)

                n_occ = len(per_occ)
                n_cols = 2
                n_rows = (n_occ + n_cols - 1) // n_cols  # ceil

                subplot_titles = [
                    f"{p['occupation']} — your n={p['client_n']}, market n={p['market_n']:,}"
                    for p in per_occ
                ]

                fig_salary_occ = make_subplots(
                    rows=n_rows, cols=n_cols,
                    subplot_titles=subplot_titles,
                    vertical_spacing=0.12,
                    horizontal_spacing=0.10,
                )

                for i, p in enumerate(per_occ):
                    row = i // n_cols + 1
                    col = i % n_cols + 1

                    # Per-occupation P02-P98 viewport clipping. Means are
                    # computed on the full distribution and stay where they
                    # are — only the histogram bars and the x-axis range
                    # honour the clipped window. Range widens to include any
                    # mean that falls outside [P02, P98] so no vline hides.
                    salaries = p['market_salaries']
                    means_present = [
                        m for m in (p['client_mean'], p['national_mean'], p['regional_mean'])
                        if not pd.isna(m)
                    ]
                    if len(salaries) > 0:
                        p02, p98 = np.percentile(salaries, [2, 98])
                        lo = min([p02] + means_present)
                        hi = max([p98] + means_present)
                        pad = (hi - lo) * 0.05 if hi > lo else max(hi * 0.05, 1)
                        x_range = [lo - pad, hi + pad]
                        hist_data = salaries[(salaries >= p02) & (salaries <= p98)]
                    else:
                        x_range = None
                        hist_data = salaries

                    fig_salary_occ.add_trace(
                        go.Histogram(
                            x=hist_data,
                            nbinsx=25,
                            marker_color=JGP_COLORS['primary'],
                            opacity=0.85,
                            showlegend=False,
                            hovertemplate='Salary: £%{x:,.0f}<br>Vacancies: %{y}<extra></extra>',
                        ),
                        row=row, col=col,
                    )

                    if not pd.isna(p['client_mean']):
                        fig_salary_occ.add_vline(
                            x=p['client_mean'], line_width=2, line_color=client_color,
                            row=row, col=col,
                        )
                    if not pd.isna(p['national_mean']):
                        fig_salary_occ.add_vline(
                            x=p['national_mean'], line_width=2, line_color=national_color,
                            row=row, col=col,
                        )
                    if not pd.isna(p['regional_mean']):
                        fig_salary_occ.add_vline(
                            x=p['regional_mean'], line_width=2, line_color=regional_color,
                            row=row, col=col,
                        )

                    if x_range is not None:
                        fig_salary_occ.update_xaxes(range=x_range, row=row, col=col)

                # Legend traces — invisible scatters drawn once on subplot (1,1)
                # so the figure-level legend has labelled rows for each line.
                fig_salary_occ.add_trace(
                    go.Scatter(
                        x=[None], y=[None], mode='lines',
                        line=dict(color=client_color, width=2),
                        name='Your mean',
                    ),
                    row=1, col=1,
                )
                fig_salary_occ.add_trace(
                    go.Scatter(
                        x=[None], y=[None], mode='lines',
                        line=dict(color=national_color, width=2),
                        name='National mean',
                    ),
                    row=1, col=1,
                )
                if any_regional:
                    fig_salary_occ.add_trace(
                        go.Scatter(
                            x=[None], y=[None], mode='lines',
                            line=dict(color=regional_color, width=2),
                            name=f"Regional mean ({client_region})",
                        ),
                        row=1, col=1,
                    )

                # Two-call layout pattern avoids Python kwarg conflicts when
                # overriding template keys (see lessons.md "Spreading
                # JGP_PLOTLY_TEMPLATE['layout']").
                fig_salary_occ.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
                fig_salary_occ.update_layout(
                    height=max(360, 240 * n_rows),
                    margin=dict(t=20, b=40, l=40, r=20),
                    showlegend=True,
                    legend=dict(
                        orientation='h',
                        yanchor='bottom', y=1.04,
                        xanchor='left', x=0,
                        font=dict(size=12),
                    ),
                    bargap=0.05,
                )
                fig_salary_occ.update_xaxes(tickformat=',', tickprefix='£')
                fig_salary_occ.update_annotations(font_size=12)

                st.plotly_chart(
                    fig_salary_occ,
                    use_container_width=True,
                    config={'displayModeBar': 'hover', 'displaylogo': False},
                )

                if not client_region:
                    st.caption(
                        "_HQ region unavailable for this client — regional benchmark "
                        "line omitted. (Common for central-government and multi-site "
                        "bodies.)_"
                    )
                elif not any_regional:
                    st.caption(
                        f"_No comparable salary data found in {client_region} for "
                        f"these occupations — regional benchmark line omitted._"
                    )

                report_figures['salary_by_occupation'] = fig_salary_occ

                # Expose for downstream commentary generator
                salary_per_occ = per_occ
                salary_client_region = client_region

    # ===================================================================
    # SECTION 08: CHANNEL PERFORMANCE (existing media performance)
    # ===================================================================
    with st.container(border=True):
        st.markdown(
            section_eyebrow('08', 'Channel performance', short='Channels')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_08_intro"]}'
            + f' <span class="info-icon" data-tooltip="{CHART_EXPLAINERS["media_performance"]}">?</span></p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )

        cat_stats = None  # Initialise before conditional block so it's in scope for PDF commentary
        if client_media is not None and len(client_media) > 0:
            # Category-level summary
            cat_stats = client_media.groupby('source_category').agg(
                total_clicks=('clicks', 'sum'),
                total_applies=('applies', 'sum'),
                vacancy_count=('entity_id_str', 'nunique')
            ).reset_index()
            cat_stats['avg_views'] = cat_stats['total_clicks'] / cat_stats['vacancy_count']
            cat_stats['avg_applies'] = cat_stats['total_applies'] / cat_stats['vacancy_count']
            cat_stats['conversion_rate'] = (cat_stats['total_applies'] / cat_stats['total_clicks'].replace(0, np.nan) * 100).fillna(0)
            cat_stats = cat_stats.sort_values('total_clicks', ascending=False)

            # On-screen: branded HTML table with inline mini-bars per row.
            max_views = cat_stats['avg_views'].max() or 1
            max_applies = cat_stats['avg_applies'].max() or 1
            BAR_BASE_PX = 140  # max bar width when value == column max

            rows_html = []
            for _, row in cat_stats.iterrows():
                views_w = max(4, int(row['avg_views'] / max_views * BAR_BASE_PX))
                applies_w = max(4, int(row['avg_applies'] / max_applies * BAR_BASE_PX))
                rows_html.append(
                    '<tr>'
                    f'<td>{row["source_category"]}</td>'
                    f'<td class="channel-num">{int(row["vacancy_count"]):,}</td>'
                    f'<td><span class="channel-bar" style="width:{views_w}px"></span>'
                    f'<span class="channel-num">{row["avg_views"]:,.1f}</span></td>'
                    f'<td><span class="channel-bar applies" style="width:{applies_w}px"></span>'
                    f'<span class="channel-num">{row["avg_applies"]:,.1f}</span></td>'
                    f'<td class="channel-pct">{row["conversion_rate"]:,.1f}%</td>'
                    '</tr>'
                )
            table_html = (
                '<table class="channel-table">'
                '<thead><tr>'
                '<th>Channel</th>'
                '<th>Vacancies</th>'
                '<th>Avg views</th>'
                '<th>Avg applies</th>'
                '<th>Conversion</th>'
                '</tr></thead>'
                f'<tbody>{"".join(rows_html)}</tbody>'
                '</table>'
            )
            st.markdown(table_html, unsafe_allow_html=True)

            # Build the bar chart silently for the PPTX export — slot in the
            # Renewals.pptx template still expects a fig.
            fig_media = go.Figure()
            fig_media.add_trace(go.Bar(
                y=cat_stats['source_category'], x=cat_stats['avg_views'],
                name='Avg views', orientation='h', marker_color=JGP_COLORS['primary']
            ))
            fig_media.add_trace(go.Bar(
                y=cat_stats['source_category'], x=cat_stats['avg_applies'],
                name='Avg applies', orientation='h', marker_color=JGP_COLORS['accent'],
                textfont=dict(color=JGP_COLORS['deep_blue']),
            ))
            fig_media.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_media.update_layout(
                barmode='group', height=max(350, len(cat_stats) * 40),
                xaxis_title="Average per vacancy", yaxis_title="",
                legend=dict(orientation='h', y=-0.15),
            )
            report_figures['media'] = fig_media

            media_commentary = generate_section_commentary('media', {
                'cat_stats': cat_stats,
                'client_name': selected_client,
            })
            st.markdown(commentary_panel(media_commentary), unsafe_allow_html=True)

            # Source-level detail — kept as a Streamlit dataframe under an expander.
            with st.expander("View by individual source", expanded=False):
                media_stats = client_media.groupby(['source_category', 'source']).agg(
                    total_clicks=('clicks', 'sum'),
                    total_applies=('applies', 'sum'),
                    vacancy_count=('entity_id_str', 'nunique')
                ).reset_index()
                media_stats['avg_views'] = media_stats['total_clicks'] / media_stats['vacancy_count']
                media_stats['avg_applies'] = media_stats['total_applies'] / media_stats['vacancy_count']
                media_stats['conversion_rate'] = (media_stats['total_applies'] / media_stats['total_clicks'].replace(0, np.nan) * 100).fillna(0)
                media_stats = media_stats.sort_values('total_clicks', ascending=False)
                detail_media = media_stats[['source_category', 'source', 'vacancy_count', 'avg_views', 'avg_applies', 'conversion_rate']].copy()
                detail_media.columns = ['Channel', 'Source', 'Vacancies', 'Avg views', 'Avg applies', 'Conversion %']
                detail_media['Avg views'] = detail_media['Avg views'].round(1)
                detail_media['Avg applies'] = detail_media['Avg applies'].round(1)
                detail_media['Conversion %'] = detail_media['Conversion %'].round(1)
                st.dataframe(detail_media, use_container_width=True, hide_index=True)
        else:
            st.info("Channel-source data isn't available yet. Build the `dashboard_media_summary` BigQuery table to enable this section.")

    # ===================================================================
    # SECTION 09: EXPORT (PowerPoint download)
    # ===================================================================
    with st.container(border=True):
        st.markdown(
            section_eyebrow('09', 'Export', short='Export')
            + f'<p class="client-section-intro">{CHART_EXPLAINERS["section_09_intro"]}</p>'
            + '<hr class="client-section-divider" />',
            unsafe_allow_html=True,
        )

        # --- Compute additional stats needed for PPTX template ---

        # Slide 2: Top-right quadrant % (vacancies above benchmark on BOTH views and applies)
        benchmarkable_df = scatter_df[scatter_df['category'] == 'Benchmarkable']
        if len(benchmarkable_df) > 0:
            top_quadrant_count = len(benchmarkable_df[
                (benchmarkable_df['views_diff_pct'] > 0) & (benchmarkable_df['applies_diff_pct'] > 0)
            ])
            top_quadrant_pct = (top_quadrant_count / len(benchmarkable_df)) * 100
        else:
            top_quadrant_pct = 0

        # Slide 2: Strongest job category (highest combined diff score, benchmarkable only)
        if len(benchmarkable_df) > 0:
            category_scores = benchmarkable_df.groupby('occupation').agg(
                combined_score=('views_diff_pct', lambda s: s.mean() + benchmarkable_df.loc[s.index, 'applies_diff_pct'].mean()),
                count=('views_diff_pct', 'count')
            ).reset_index()
            # Need at least 2 vacancies in occupation for the category to be considered "strongest"
            category_scores = category_scores[category_scores['count'] >= 2]
            if len(category_scores) > 0:
                top_category = category_scores.sort_values('combined_score', ascending=False).iloc[0]['occupation']
            else:
                top_category = benchmarkable_df.iloc[0]['occupation'] if len(benchmarkable_df) > 0 else 'N/A'
        else:
            top_category = 'N/A'

        # Slide 5: Build the new charts (only if spend entered)
        rate_card_total_val = rate_card_price * num_jobs if annual_spend > 0 else 0
        saving_pct_val = ((rate_card_total_val - annual_spend) / rate_card_total_val * 100) if rate_card_total_val > 0 else 0

        if annual_spend > 0:
            # Stacked bar: Your Spend (bottom) + Saving (top) = Rate Card Total
            saving_amount = max(rate_card_total_val - annual_spend, 0)
            fig_spend_stack = go.Figure()
            fig_spend_stack.add_trace(go.Bar(
                x=['Total Value'],
                y=[annual_spend],
                name='Your Spend',
                marker_color=JGP_COLORS['primary'],
                text=[f"£{annual_spend:,.0f}"],
                textposition='inside',
                textfont=dict(color='white', size=14),
                width=0.7,
            ))
            fig_spend_stack.add_trace(go.Bar(
                x=['Total Value'],
                y=[saving_amount],
                name='Saving vs Rate Card',
                marker_color=JGP_COLORS['accent'],
                text=[f"£{saving_amount:,.0f}"],
                textposition='inside',
                textfont=dict(color=JGP_COLORS['deep_blue'], size=14),
                width=0.7,
            ))
            fig_spend_stack.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
            fig_spend_stack.update_layout(
                barmode='stack',
                title=f"Your Spend vs Rate Card Value (Saving: {saving_pct_val:.0f}%)",
                yaxis_title="GBP",
                height=400,
                bargap=0.05,
                showlegend=True,
                plot_bgcolor='rgba(0,0,0,0)',
                legend=dict(orientation='h', y=-0.15),
            )
            report_figures['spend_vs_ratecard'] = fig_spend_stack

        # Cost per apply by occupation chart (always built when spend > 0)
        cpa_by_occ_fig = None
        roi_by_type_full = None
        if annual_spend > 0:
            roi_by_type_full = client_df.groupby('occupation').agg(
                total_applies=('applies', 'sum'),
                job_count=('clicks', 'count')
            ).reset_index()
            roi_by_type_full = roi_by_type_full[roi_by_type_full['total_applies'] > 0]
            if len(roi_by_type_full) > 0:
                roi_by_type_full['cost_allocated'] = annual_spend * (roi_by_type_full['job_count'] / roi_by_type_full['job_count'].sum())
                roi_by_type_full['cost_per_apply'] = roi_by_type_full['cost_allocated'] / roi_by_type_full['total_applies']
                roi_by_type_full = roi_by_type_full.sort_values('cost_per_apply', ascending=True)

                cpa_by_occ_fig = go.Figure()
                cpa_by_occ_fig.add_trace(go.Bar(
                    y=roi_by_type_full['occupation'],
                    x=roi_by_type_full['cost_per_apply'],
                    orientation='h',
                    marker_color=JGP_COLORS['blue'],
                    text=roi_by_type_full['cost_per_apply'].apply(lambda x: f"£{x:,.2f}"),
                    textposition='outside',
                    textfont=dict(color=JGP_COLORS['deep_blue']),
                ))
                cpa_by_occ_fig.update_layout(**JGP_PLOTLY_TEMPLATE['layout'])
                cpa_by_occ_fig.update_layout(
                    title="Cost per Apply by Occupation",
                    height=max(350, len(roi_by_type_full) * 32),
                    xaxis_title="Cost per Apply (GBP)",
                    yaxis_title="",
                    plot_bgcolor='rgba(0,0,0,0)',
                )
                report_figures['cost_per_app_by_occupation'] = cpa_by_occ_fig

        # --- Build structured commentary for PPTX template ---
        scatter_struct = generate_section_commentary_structured('benchmark_scatter', {
            'total_count': len(scatter_df),
            'benchmarkable_count': len(scatter_df[scatter_df['category'] == 'Benchmarkable']),
            'zero_applies_count': len(scatter_df[scatter_df['category'].str.startswith('Zero')]),
            'no_benchmark_count': len(scatter_df[scatter_df['category'] == 'Low Sample (No Benchmark)']),
            'top_performers': top_performers,
            'client_name': selected_client,
        })
        average_struct = generate_section_commentary_structured('benchmark_average', {
            'client_avg_clicks': client_avg_clicks,
            'benchmark_avg_clicks': benchmark_avg_clicks,
            'client_avg_applies': client_avg_applies,
            'benchmark_avg_applies': benchmark_avg_applies,
            'num_jobs': len(client_df),
            'client_name': selected_client,
        })
        postings_struct = generate_section_commentary_structured('postings', {
            'total_jobs': len(client_df),
            'total_applies': int(client_df['applies'].sum()),
            'by_type': by_type,
            'client_name': selected_client,
        })
        roi_struct = generate_section_commentary_structured('roi', {
            'annual_spend': annual_spend,
            'rate_card_total': rate_card_total_val,
            'num_jobs': num_jobs,
            'cost_per_apply': annual_spend / total_applies_val if total_applies_val > 0 else 0,
            'saving_pct': saving_pct_val,
            'roi_by_type': roi_by_type_full if (annual_spend > 0 and roi_by_type_full is not None and len(roi_by_type_full) > 0) else None,
            'client_name': selected_client,
        })
        media_struct = generate_section_commentary_structured('media', {
            'cat_stats': cat_stats,
            'client_name': selected_client,
        })
        salary_struct = generate_section_commentary_structured('salary', {
            'per_occ': salary_per_occ,
            'client_name': selected_client,
            'client_region': salary_client_region,
        })

        # --- Build report_metrics dict (matches template tag names) ---
        report_metrics = {
            # Slide 1
            'client_name': selected_client,
            'PERIOD_START': str(report_start),
            'PERIOD_END': str(report_end),

            # Slide 2 stats
            'stat_total_jobs': f"{num_jobs:,}",
            'stat_top_quadrant_pct': f"{top_quadrant_pct:.0f}",
            'stat_top_category': top_category,

            # Slide 2 commentary
            'commentary_benchmark_intro': scatter_struct['intro'],
            'commentary_benchmark_point_1': scatter_struct['point_1'],
            'commentary_benchmark_point_2': scatter_struct['point_2'],
            'commentary_benchmark_point_3': scatter_struct['point_3'],

            # Slide 3 stats
            'stat_benchmark_average_views': f"{benchmark_avg_clicks:,.0f}",
            'stat_your_jobs_average_views': f"{client_avg_clicks:,.0f}",
            'stat_benchmark_average_applies': f"{benchmark_avg_applies:,.1f}",
            'stat_your_jobs_average_applies': f"{client_avg_applies:,.1f}",

            # Slide 3 commentary
            'commentary_average_intro': average_struct['intro'],
            'commentary_average_point_1': average_struct['point_1'],
            'commentary_average_point_2': average_struct['point_2'],

            # Slide 4 commentary
            'commentary_postings_intro': postings_struct['intro'],
            'commentary_postings_point_1': postings_struct['point_1'],
            'commentary_postings_point_2': postings_struct['point_2'],

            # Slide 5 stats (ROI)
            'stat_cost_per_job': f"£{annual_spend / num_jobs:,.2f}" if (annual_spend > 0 and num_jobs > 0) else "—",
            'stat_cost_per_view': f"£{annual_spend / total_clicks:,.2f}" if (annual_spend > 0 and total_clicks > 0) else "—",
            'stat_cost_per_apply': f"£{annual_spend / total_applies_val:,.2f}" if (annual_spend > 0 and total_applies_val > 0) else "—",

            # Slide 5 commentary
            'commentary_roi_intro': roi_struct['intro'],
            'commentary_roi_point_1': roi_struct['point_1'],
            'commentary_roi_point_2': roi_struct['point_2'],

            # Slide 6 commentary (Salary Benchmark)
            'commentary_salary_intro': salary_struct['intro'],
            'commentary_salary_point_1': salary_struct['point_1'],
            'commentary_salary_point_2': salary_struct['point_2'],
            'commentary_salary_point_3': salary_struct['point_3'],

            # Slide 7 commentary
            'commentary_media_intro': media_struct['intro'],
            'commentary_media_point_1': media_struct['point_1'],
            'commentary_media_point_2': media_struct['point_2'],
            'commentary_media_point_3': media_struct['point_3'],

            # Static chart explainers — sourced from module-level CHART_EXPLAINERS
            # so on-screen captions and PPTX placeholders share one source of truth.
            **{f'chart_explainer_{k}': v for k, v in CHART_EXPLAINERS.items()},

            # Slide 7 contact
            'contact_name': contact_name or 'Your Account Manager',
            'contact_title': contact_title or 'Account Director',
            'contact_email': contact_email or 'team@jobsgopublic.com',
            'contact_phone': contact_phone or '020 7427 8250',
        }

        # --- Map template chart tags to figures ---
        pptx_figures = {
            'benchmark_scatter': report_figures.get('scatter'),
            'benchmark_average': report_figures.get('benchmark_combined'),
            'postings_by_type': report_figures.get('postings'),
            'spend_vs_ratecard': report_figures.get('spend_vs_ratecard'),
            'cost_per_app_by_occupation': report_figures.get('cost_per_app_by_occupation'),
            'salary_by_occupation': report_figures.get('salary_by_occupation'),
            'media_performance': report_figures.get('media'),
        }

        # --- Export CTA panel + download button ---
        st.markdown(
            export_cta_panel(
                heading="Export this report",
                lede=(
                    f"Download a branded PowerPoint version of the {selected_client} "
                    "report to share with the client or drop into a renewals deck. "
                    "All charts, KPIs and commentary are included; tweak any text in "
                    "PowerPoint, then File → Export → PDF for the final copy."
                ),
            ),
            unsafe_allow_html=True,
        )

        template_path = 'Renewals.pptx'
        try:
            pptx_bytes = generate_client_report_pptx(report_metrics, pptx_figures, template_path)
            st.download_button(
                "Download PowerPoint report",
                data=pptx_bytes,
                file_name=f"advertising_report_{selected_client.replace(' ', '_')}_{report_start}_{report_end}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary"
            )
        except FileNotFoundError:
            st.error(f"Template file not found at `{template_path}`. Make sure `Renewals.pptx` is in the project root.")
        except Exception as e:
            st.warning("PowerPoint generation requires `python-pptx` and `kaleido`.")
            st.caption(f"Error: {type(e).__name__}: {e}")


def generate_client_report_pptx(metrics, figures, template_path):
    """Generate a PowerPoint report by filling a branded template.

    Args:
        metrics: dict mapping placeholder names (without {{}}) to string values
        figures: dict mapping chart slot names (e.g. 'benchmark_scatter') to plotly figures
        template_path: path to the .pptx template file

    Returns: bytes of the populated .pptx file.

    Replaces text placeholders like {{tag_name}} with metrics[tag_name].
    Replaces chart placeholders like {{chart:slot_name}} with PNG images of figures[slot_name].
    """
    prs = Presentation(template_path)

    # --- Helper: render a Plotly figure to PNG bytes at a target aspect ratio ---
    def _fig_to_png(fig, slot_width_emu=None, slot_height_emu=None):
        """Render a Plotly figure to a high-DPI white-background PNG.

        Critical sizing rule: the logical canvas is sized to the placeholder's
        displayed pixel dimensions (1 px ≈ 9525 EMU at 96 DPI). Plotly font
        sizes are in pixels *relative to the canvas*, so when the PNG is fitted
        back into the placeholder the displayed font matches what we set.
        The previous fixed 1800 px canvas made fonts shrink to 4–6 pt on small
        placeholders like spend_vs_ratecard, regardless of how big the source
        font was set.

        Crispness comes from `scale=6` — a pure DPI multiplier on the rendered
        bitmap that keeps the logical canvas (and therefore font sizing) intact.
        Output is 6× the canvas dimensions, e.g. an 800 × 500 logical canvas
        renders to a 4800 × 3000 PNG.

        Font sizes are floored at 14 px ≈ 10.5 pt — the "no smaller than 10 pt"
        readability target. Axis titles, bar value labels and chart titles step
        up from there.
        """
        if fig is None:
            return None
        try:
            # 1 EMU = 1/914400 inch; ~9525 EMU per displayed pixel at 96 DPI.
            EMU_PER_PX = 9525
            # Small floor only as a safeguard against Plotly choking on tiny
            # canvases — well below any realistic slide placeholder size.
            MIN_W, MIN_H = 320, 220
            if slot_width_emu and slot_height_emu:
                width = max(int(slot_width_emu / EMU_PER_PX), MIN_W)
                height = max(int(slot_height_emu / EMU_PER_PX), MIN_H)
            else:
                width, height = 1280, 720

            fig_export = go.Figure(fig.to_dict())
            # Standard text scheme — DM Sans Regular (weight 400). Sizes match
            # the last pass: with canvas = placeholder px, source px ≈ displayed
            # px ≈ pt × 1.333:
            #   13 px ≈ 10 pt → tick labels (x AND y), legend, scatter labels
            #   16 px ≈ 12 pt → axis titles
            #   12 px ≈  9 pt → bar/data value labels (explicit floor)
            #   13 px ≈ 10 pt → in-chart annotations
            #
            # Background is fully transparent so the slide colour/shapes show
            # through the chart PNG (PNG's alpha channel preserved by kaleido).
            DM_SANS = 'DM Sans, sans-serif'
            TRANSPARENT = 'rgba(0,0,0,0)'
            fig_export.update_layout(
                paper_bgcolor=TRANSPARENT,
                plot_bgcolor=TRANSPARENT,
                font=dict(family=DM_SANS,
                          color=JGP_COLORS['deep_blue'], size=13),
                # Drop the chart title — every slide has its own hard-coded title.
                title=dict(text=''),
                # Top margin tightened so the plot area extends as close to the
                # top of the slot as possible.
                margin=dict(l=80, r=40, t=15, b=70, pad=10),
                legend=dict(
                    # Legend stays partially translucent white so points behind
                    # legend text remain readable on coloured slide backgrounds.
                    bgcolor='rgba(255,255,255,0.6)',
                    font=dict(family=DM_SANS, size=13),
                ),
            )
            fig_export.update_xaxes(
                automargin=True,
                title_font=dict(family=DM_SANS, size=16),
                tickfont=dict(family=DM_SANS, size=13),
            )
            fig_export.update_yaxes(
                automargin=True,
                title_font=dict(family=DM_SANS, size=16),
                tickfont=dict(family=DM_SANS, size=13),
            )

            # Trace-level overrides for export visibility. Disable Plotly's
            # constraintext (defaults to 'both', which silently shrinks bar
            # value labels to fit inside narrow bars — the root cause of the
            # tiny "£spend" labels on spend_vs_ratecard).
            for trace in fig_export.data:
                t = trace.type
                if t == 'bar':
                    preserved_color = None
                    if trace.textfont is not None and trace.textfont.color is not None:
                        preserved_color = trace.textfont.color
                    # 12 px ≈ 9 pt — explicit per user request.
                    tf = {'family': DM_SANS, 'size': 12}
                    if preserved_color:
                        tf['color'] = preserved_color
                    trace.textfont = tf
                    trace.constraintext = 'none'
                elif t == 'scatter':
                    # Bump dot markers and line widths for slide visibility.
                    marker = trace.marker
                    if marker is not None:
                        cur_size = marker.size
                        if isinstance(cur_size, (int, float)):
                            marker.size = max(cur_size, 14)
                        elif cur_size is None:
                            marker.size = 14
                    line = trace.line
                    if line is not None:
                        cur_w = line.width
                        if isinstance(cur_w, (int, float)):
                            line.width = max(cur_w, 3.5)
                    # Match scatter data labels to the universal 10 pt standard.
                    tf_existing = trace.textfont
                    if tf_existing is not None:
                        trace.textfont = {
                            'family': DM_SANS,
                            'size': 13,
                            'color': (tf_existing.color
                                      if tf_existing.color
                                      else JGP_COLORS['deep_blue']),
                        }

            # Annotation font (e.g. benchmark line callout on benchmark_average).
            if fig_export.layout.annotations:
                for ann in fig_export.layout.annotations:
                    cur = (ann.font.size if ann.font and ann.font.size else None)
                    ann.font = dict(
                        family=DM_SANS,
                        size=max(cur or 13, 13),
                        color=(ann.font.color if ann.font and ann.font.color
                               else JGP_COLORS['deep_blue']),
                    )

            # scale=6 → PNG is 6x the logical canvas (e.g. 800×500 → 4800×3000).
            # Maximum crispness at the cost of a few seconds per chart and
            # ~1–3 MB per PNG; acceptable for client renewal reports.
            return fig_export.to_image(format='png', width=width, height=height, scale=6)
        except Exception:
            return None

    # --- Helper: replace text in a single shape's text frame, preserving formatting ---
    def _replace_text_in_shape(shape, replacements):
        """Walk runs/paragraphs and replace {{tag}} occurrences. Handles tags split across runs."""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            # First try simple per-run replacement (works when tag is fully in one run)
            for run in paragraph.runs:
                text = run.text
                if '{{' in text:
                    for tag, val in replacements.items():
                        placeholder = '{{' + tag + '}}'
                        if placeholder in text:
                            text = text.replace(placeholder, str(val))
                    run.text = text

            # Fallback: if a tag spans multiple runs, the above won't catch it.
            # Concatenate paragraph text, replace, and put it all in the first run.
            full_text = ''.join(r.text for r in paragraph.runs)
            if '{{' in full_text and any('{{' + tag + '}}' in full_text for tag in replacements):
                new_text = full_text
                for tag, val in replacements.items():
                    new_text = new_text.replace('{{' + tag + '}}', str(val))
                if new_text != full_text:
                    if paragraph.runs:
                        paragraph.runs[0].text = new_text
                        for r in paragraph.runs[1:]:
                            r.text = ''

    # --- Step 1: Find all chart placeholders, capture position/size, queue for replacement ---
    # We do this before text replacement so we can find {{chart:xxx}} markers.
    chart_replacements = []  # list of (slide, shape, slot_name)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                m = re.search(r'\{\{chart:([^}]+)\}\}', text)
                if m:
                    chart_replacements.append((slide, shape, m.group(1)))

    # --- Step 2: Replace text placeholders on every shape across all slides ---
    # `slide_number` and `total_slides` are computed per-slide so a single
    # template footer like `{{slide_number}} / {{total_slides}}` renders as
    # "1 / 9" on slide 1, "2 / 9" on slide 2, etc.
    text_replacements = {k: v for k, v in metrics.items()}
    total_slides = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        per_slide = {
            **text_replacements,
            'slide_number': idx,
            'total_slides': total_slides,
        }
        for shape in slide.shapes:
            _replace_text_in_shape(shape, per_slide)

    # --- Step 3: Replace chart placeholders with images ---
    for slide, shape, slot_name in chart_replacements:
        fig = figures.get(slot_name)
        # Capture original geometry before deleting the shape
        left, top, width, height = shape.left, shape.top, shape.width, shape.height

        # Remove the placeholder shape
        sp = shape._element
        sp.getparent().remove(sp)

        if fig is not None:
            png_bytes = _fig_to_png(fig, slot_width_emu=width, slot_height_emu=height)
            if png_bytes:
                slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width=width, height=height)
            else:
                # Failed to render — leave a small note in place
                txt_box = slide.shapes.add_textbox(left, top, width, height)
                txt_box.text_frame.text = '[Chart unavailable]'
        # If fig is None, just remove the placeholder silently

    # --- Step 4: Clean up any remaining unreplaced {{tags}} (set to empty so they don't show) ---
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = ''.join(r.text for r in paragraph.runs)
                    if re.search(r'\{\{[^}]+\}\}', full_text):
                        cleaned = re.sub(r'\{\{[^}]+\}\}', '', full_text)
                        if paragraph.runs:
                            paragraph.runs[0].text = cleaned
                            for r in paragraph.runs[1:]:
                                r.text = ''

    # --- Step 5: Output to bytes ---
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
